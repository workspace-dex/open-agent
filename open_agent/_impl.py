#!/usr/bin/env python3
"""
Open-Agent — local AI agent for consumer hardware.
One file. No cloud. Yours.
"""

# ══════════════════════════════════════════════════════════════════════════
#  IMPORTS
# ══════════════════════════════════════════════════════════════════════════

import argparse
import asyncio
import feedparser
import json
import os
import re
import select
import shutil
import subprocess
import sys
import termios
import threading
import tty
import uuid
import sqlite3
from datetime import datetime
from pathlib import Path
from typing import Optional, Tuple, List, Set, Dict

# Custom AgentLoop for parallel execution with real-time streaming
from open_agent.agent_loop import AgentLoop as CustomAgentLoop, TOOL_SCHEMAS

# Real-time streaming flag - set when text is being streamed live
_was_streamed_live = [False]

# ══════════════════════════════════════════════════════════════════════════
#  SIGINT / Ctrl+C  (mirrors hermes-agent tools/interrupt.py pattern)
# ══════════════════════════════════════════════════════════════════════════
# Thread-safe: threading.Event is used so the signal handler (runs in the
# main thread) can safely flag the agent thread to interrupt.

def _sigint_handler(signum, frame):
    """Called on Ctrl+C — signal the agent loop to abort at the next safe point."""
    from open_agent.agent_loop import interrupt_agent
    interrupt_agent()
    # Print immediately (stdout may be buffered by Rich/Live)
    try:
        sys.stdout.write("\n\033[2m  ⚡  interrupted — finishing current tool…\033[0m\n")
        sys.stdout.flush()
    except OSError:
        pass
    # Re-raise KeyboardInterrupt so the main loop's try/except catches it
    raise KeyboardInterrupt()

try:
    import signal
    signal.signal(signal.SIGINT, _sigint_handler)
    if hasattr(signal, 'SIGHUP'):
        signal.signal(signal.SIGHUP, _sigint_handler)
except (ValueError, OSError):
    # SIGHUP doesn't exist on Windows; SIGINT may already be handled
    pass

def _set_streamed():
    """Mark that text was already streamed live - don't print again."""
    _was_streamed_live[0] = True


# ══════════════════════════════════════════════════════════════════════════
#  ERROR CLASSIFICATION  (hermes-agent pattern)
# ══════════════════════════════════════════════════════════════════════════

# Each entry: (regex, hint, retryable, compressible, rotate)
_FAILCLASS_RE = [
    (re.compile(r'\bPermissionError\b'),               "Check file permissions: chmod +x, chown, or run as correct user.",                     True,  False, False),
    (re.compile(r'\bFileNotFoundError\b'),              "File doesn't exist. Check for typos (Linux is case-sensitive).",                   True,  False, False),
    (re.compile(r'\bIsADirectoryError\b'),              "Expected a file, got a directory. Check your path.",                                True,  False, False),
    (re.compile(r'\bNotADirectoryError\b'),             "Expected a directory, got a file. Check your path.",                               True,  False, False),
    (re.compile(r'\bConnectionRefusedError\b'),           "Server not running. Start it first or check the port.",                             True,  False, False),
    (re.compile(r'\bConnectionError\b'),                "Network error. Check URL, firewall, or internet connection.",                         True,  False, False),
    (re.compile(r'\bTimeoutExpired\b'),                 "Operation timed out. Break into smaller steps or increase timeout.",                 True,  False, False),
    (re.compile(r'\bBadStatusLine\b|\bRemoteProtocolError\b'), "Protocol error from server. Retry — server may have restarted.",            True,  False, False),
    (re.compile(r'\bJSONDecodeError\b|\bjson\..JSONDecodeError\b'), "Invalid JSON. Check syntax: missing quotes, commas, or brackets.",        True,  False, False),
    (re.compile(r'\bSyntaxError\b'),                   "Python syntax error. Check parentheses, colons, indentation.",                         True,  False, False),
    (re.compile(r'\bIndentationError\b'),               "Indentation error. Use 4 spaces (not tabs). Fix alignment of function/class body.",  True,  False, False),
    (re.compile(r'\bNameError\b'),                     "Variable not defined. Define it before use: x = 'value'",                           True,  False, False),
    (re.compile(r'\bTypeError\b'),                     "Type mismatch. Check you're passing the right types (str vs int, etc.)",             True,  False, False),
    (re.compile(r'\bValueError\b'),                     "Invalid value. Check argument types and valid ranges.",                               True,  False, False),
    (re.compile(r'\bAttributeError\b'),                 "Object has no attribute. Check the API docs or use dir(obj).",                     True,  False, False),
    (re.compile(r'\bKeyError\b'),                       "Key not found in dict. Check keys with dict.keys() or use .get().",                 True,  False, False),
    (re.compile(r'\bIndexError\b'),                     "Index out of range. Check list length before accessing.",                            True,  False, False),
    (re.compile(r'\bImportError\b'),                     "Module not installed. Run: pip install <module>",                                    True,  False, False),
    (re.compile(r'\bModuleNotFoundError\b'),             "Install with: pip install <package-name>",                                          True,  False, False),
    (re.compile(r'\bNoSuchCommand\b|\bCommandNotFoundError\b'), "Command not found. Install it or check your PATH.",                         True,  False, False),
    (re.compile(r'\b403\b|\bforbidden\b', re.I),        "Access denied (403). Check API key permissions or rate limits.",                   True,  True,  False),
    (re.compile(r'\b404\b|\bnot found\b', re.I),        "Resource not found (404). Check URL or resource ID.",                               True,  False, False),
    (re.compile(r'\b429\b|\brate limit\b', re.I),        "Rate limited (429). Wait and retry. Consider caching repeated requests.",            True,  True,  True),
    (re.compile(r'\b500\b|\bInternal Server Error\b'),   "Server error (500). Retry in a few seconds — likely transient.",                   True,  True,  False),
    (re.compile(r'\b502\b|\b503\b|\b504\b'),             "Gateway error from server. Retry — usually transient.",                              True,  True,  False),
]

def classify_error(error_msg: str) -> dict:
    """
    Structured error classification. Returns {hint, retryable, compressible, rotate}.
    - retryable: worth retrying the same operation
    - compressible: retry with shorter context
    - rotate: retry with different model/provider
    """
    for pattern, hint, retryable, compressible, rotate in _FAILCLASS_RE:
        if pattern.search(error_msg):
            return dict(hint=hint, retryable=retryable, compressible=compressible, rotate=rotate)
    return dict(hint="", retryable=True, compressible=False, rotate=False)

def _analyze_error(error_msg: str) -> str:
    """Single-hint version of classify_error."""
    cls = classify_error(error_msg)
    return f"Hint: {cls['hint']}" if cls['hint'] else ""


# ══════════════════════════════════════════════════════════════════════════
#  JITTERED EXPONENTIAL BACKOFF
# ══════════════════════════════════════════════════════════════════════════

def jittered_backoff(attempt: int, base: float = 0.5, max_delay: float = 60.0) -> float:
    """
    Exponential backoff with uniform jitter.
    Exponential base is capped BEFORE jitter to prevent wild values.
    Delay = min(base * 2^attempt, max_delay) + random(0, 0.5 * capped_base).
    Thread-safe, no shared state.
    """
    import random as _rand
    exp = min(base * (2 ** attempt), max_delay)
    jitter = _rand.uniform(0, exp * 0.5) if attempt > 0 else 0
    return min(exp + jitter, max_delay)


# ══════════════════════════════════════════════════════════════════════════
#  TOOL REGISTRY  — structured metadata + dynamic lifecycle
# ══════════════════════════════════════════════════════════════════════════
from dataclasses import dataclass, field
from typing import Any, Callable

@dataclass(frozen=True)
class ToolSpec:
    """
    Structured metadata for a single tool.

    Fields:
        name        — canonical identifier (matches function name)
        description — one-line purpose for the model
        danger      — 0=safe, 1=moderate (writes/commands), 2=high (system)
        examples    — usage examples shown to the model
        categories  — tag the tool for filtering/enabling
        enabled     — runtime toggle (can disable per session/task)
        validator   — optional (args → None or raises ValueError)
    """
    name:        str
    description: str
    danger:      int = 0
    examples:    tuple[str, ...] = field(default_factory=tuple)
    categories:  tuple[str, ...] = field(default_factory=tuple)
    enabled:     bool = True
    validator:   Callable[[dict], None] | None = None

    def disable(self) -> None:
        object.__setattr__(self, "enabled", False)

    def enable(self) -> None:
        object.__setattr__(self, "enabled", True)


class ToolRegistry:
    """
    Centralised tool registry — single source of truth for all tool metadata.

    Operations:
        register(spec)        — add a tool
        get(name)             — get spec or None
        list_all(enabled_only) — all specs
        list_by_category(cat) — specs filtered by category
        enable(name)          — runtime enable
        disable(name)         — runtime disable
        validate(name, args)  — run validator, raise ValueError on bad args
    """

    def __init__(self) -> None:
        self._specs: dict[str, ToolSpec] = {}

    def register(self, spec: ToolSpec) -> None:
        self._specs[spec.name] = spec

    def get(self, name: str) -> ToolSpec | None:
        return self._specs.get(name)

    def list_all(self, enabled_only: bool = True) -> list[ToolSpec]:
        specs = self._specs.values()
        if enabled_only:
            specs = [s for s in specs if s.enabled]
        return sorted(specs, key=lambda s: s.name)

    def list_by_category(self, category: str) -> list[ToolSpec]:
        return [s for s in self._specs.values() if s.enabled and category in s.categories]

    def enable(self, name: str) -> bool:
        spec = self._specs.get(name)
        if spec:
            spec.enable()
            return True
        return False

    def disable(self, name: str) -> bool:
        spec = self._specs.get(name)
        if spec:
            spec.disable()
            return True
        return False

    def validate(self, name: str, args: dict) -> None:
        """Run the tool's validator if present; no-op otherwise."""
        spec = self._specs.get(name)
        if spec and spec.validator:
            spec.validator(args)

    def to_prompt_lines(self) -> str:
        """Render all enabled tools as description lines for the system prompt."""
        lines = []
        for spec in self.list_all():
            ex_lines = ""
            if spec.examples:
                ex_lines = "\n  Examples: " + "; ".join(spec.examples)
            lines.append(
                f"{spec.name}({spec.description}){ex_lines}"
            )
        return "\n".join(lines)


# ══════════════════════════════════════════════════════════════════════════
#  TOOLSET SYSTEM (Lazy-loading like hermes-agent)
# ══════════════════════════════════════════════════════════════════════════

# Minimal toolset definitions - only load what you need
# These match the ACTUAL @agent.tool functions defined in this file
TOOLSETS = {
    "web": {
        "description": "Web search and content extraction",
        "tools": ["web_search", "cached_web_search", "smart_research", "fetch_page"],
        "keywords": ["search", "find", "google", "latest", "news", "web", "url", "download", "fetch", "scrape"],
    },
    "file": {
        "description": "File read, write, and search operations", 
        "tools": ["read_file", "write_file"],
        "keywords": ["file", "read", "write", "save", "create", "edit", "open", "view"],
    },
    "terminal": {
        "description": "Shell command execution",
        "tools": ["run_terminal"],
        "keywords": ["run", "execute", "command", "shell", "bash", "python", "install", "git", "npm"],
    },
    "memory": {
        "description": "Persistent memory and session search",
        "tools": ["update_memory", "read_memory", "search_sessions", "update_user_profile"],
        "keywords": ["remember", "memory", "history", "session", "past", "preference", "profile"],
    },
    "obsidian": {
        "description": "Obsidian vault notes",
        "tools": ["search_obsidian", "read_obsidian_note", "write_obsidian_note"],
        "keywords": ["obsidian", "note", "vault", "markdown"],
    },
    "office": {
        "description": "Office document creation",
        "tools": ["create_pptx"],
        "keywords": ["ppt", "presentation", "slides", "deck", "powerpoint"],
    },
    "soul": {
        "description": "Load extended instructions",
        "tools": ["load_soul"],
        "keywords": ["soul", "instructions", "extended", "behavior"],
    },
}

# Keywords that trigger ALL toolsets (fallback)
_FALLBACK_KEYWORDS = {"all", "everything", "full", "complete"}


def _detect_toolset(user_input: str) -> Set[str]:
    """Detect which toolsets to load based on user input."""
    user_lower = user_input.lower()
    detected = set()
    
    for toolset_name, config in TOOLSETS.items():
        for keyword in config.get("keywords", []):
            if keyword in user_lower:
                detected.add(toolset_name)
                break
    
    # If nothing detected, load minimal set (web + terminal)
    if not detected:
        detected = {"web", "terminal"}
    
    return detected


def get_toolset_tools(toolset_names: Set[str]) -> List[str]:
    """Get list of tool names for given toolsets."""
    tools = []
    for name in toolset_names:
        if name in TOOLSETS:
            tools.extend(TOOLSETS[name]["tools"])
    return tools

import httpx
from pydantic_ai import Agent, RunContext
from pydantic_ai.messages import ModelMessage
from pydantic_ai.models.openai import OpenAIChatModel
from pydantic_ai.providers.openai import OpenAIProvider

# ── Enhanced tools integration (import without modifying original) ────────────
try:
    from open_agent.enhanced import ENABLED
    _ENHANCED = True
except ImportError:
    _ENHANCED = False
    ENABLED = {}

# Thinking logger (hermes-style logs)
_thinking_chain = None
_thinking_original_td = None

def _init_thinking_logger():
    """Initialize thinking logger on first use."""
    global _thinking_chain, _thinking_original_td
    if _thinking_chain is None and _ENHANCED:
        try:
            from open_agent.thinking_logger import ThinkingChain, ToolBadgeLogger
            _thinking_chain = ThinkingChain(verbose=True)
            if hasattr(sys, '_ ThinkingChain'):
                pass
            _w(f"{DIM}  [Thinking logger active]{RESET}\n")
        except ImportError:
            pass

def _patch_td_with_thinking():
    """Patch td (ToolBadge) with thinking logger."""
    global _thinking_chain, _thinking_original_td
    if _ENHANCED and _thinking_chain is None:
        try:
            from open_agent.thinking_logger import ThinkingChain, ToolBadgeLogger
            _thinking_chain = ThinkingChain(verbose=True)
            _thinking_original_td = td
            # Create wrapper
            class _LoggedBadge:
                def __init__(self, orig, chain):
                    self._orig = orig
                    self._chain = chain
                def start(self, name, detail=""):
                    self._chain.tool_call(name, {"detail": str(detail)[:120]})
                    self._orig.start(name, detail)
                def done(self, name, summary="", ok=True, chars=0):
                    self._chain.tool_result(name, ok, chars, str(summary)[:80])
                    self._orig.done(name, summary, ok, chars)
                def info(self, label=""):
                    self._chain.context_info(label)
                    self._orig.info(label)
            globals()['td'] = _LoggedBadge(td, _thinking_chain)
        except ImportError:
            pass

def _unpatch_td():
    """Restore original ToolBadge."""
    global _thinking_original_td
    if _thinking_original_td is not None:
        globals()['td'] = _thinking_original_td
        _thinking_original_td = None

from rich import box
from rich.console import Console
from rich.markdown import Markdown
from rich.panel import Panel
from rich.syntax import Syntax
from rich.table import Table
from rich.text import Text
from rich.theme import Theme
from rich.align import Align
from rich.columns import Columns
from rich.live import Live


# ══════════════════════════════════════════════════════════════════════════
#  ANSI HELPERS  (direct stdout — no Rich buffering during streaming)
# ══════════════════════════════════════════════════════════════════════════

ESC        = "\x1b"
RESET      = f"{ESC}[0m"
BOLD       = f"{ESC}[1m"
DIM        = f"{ESC}[2m"
C_CYAN     = f"{ESC}[38;5;117m"
C_GREEN    = f"{ESC}[38;5;114m"
C_AMBER    = f"{ESC}[38;5;214m"
C_PURPLE   = f"{ESC}[38;5;183m"
C_TEAL     = f"{ESC}[38;5;115m"
C_RED      = f"{ESC}[38;5;203m"
C_SLATE    = f"{ESC}[38;5;60m"
C_WHITE    = f"{ESC}[38;5;252m"
C_BLUE     = f"{ESC}[38;5;75m"
CLEAR_LINE = f"\r{ESC}[2K"
CLEAR_EOL  = f"{ESC}[K"        # clear from cursor to end of line
CLEAR_SCR  = f"{ESC}[2J{ESC}[H" # clear screen + home

# ── Color support detection ─────────────────────────────────────────────
_ANSI_RE       = re.compile(r"\x1b\[[0-9;]*[mK]")
_use_ansi_cache: bool | None = None

def should_use_color() -> bool:
    """True when color output is supported and not disabled."""
    global _use_ansi_cache
    if _use_ansi_cache is None:
        _use_ansi_cache = (
            sys.stdout.isatty()
            and os.getenv("NO_COLOR") is None
            and os.getenv("TERM") not in ("dumb", "")
            and os.getenv("TERM") is not None
        )
    return _use_ansi_cache

def _w(s: str):
    """Write to stdout — broken-pipe safe; strips ANSI when color disabled."""
    if not should_use_color():
        s = _ANSI_RE.sub("", s)
    try:
        sys.stdout.write(s)
        sys.stdout.flush()
    except OSError:
        pass  # stdout closed — nothing to do

def _ln(s: str = ""):
    _w(s + "\n")


def _normalize_output(text: str) -> str:
    """Tighten assistant output before rendering and saving. Strip ANSI sequences."""
    text = _ANSI_RE.sub("", text)          # remove color codes
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = text.replace("", "")
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def truncate_result(content: str, max_chars: int = None) -> str:
    """Truncate large tool results for token efficiency."""
    max_chars = max_chars or CFG.get("max_tool_result_chars", 6_000)
    if len(content) <= max_chars:
        return content
    # Find a good break point (last newline within limit)
    truncated = content[:max_chars]
    last_nl = truncated.rfind("\n")
    if last_nl > max_chars // 2:
        truncated = truncated[:last_nl + 1]
    return truncated + f"\n\n[Truncated: {len(content):,} chars → {max_chars:,}]"


def _looks_like_planning(text: str) -> bool:
    lowered = text.lower()
    planning_markers = (
        "i will", "i'll", "let me", "i am going to", "i’m going to",
        "i will now", "first i will", "i’ll now", "i will fetch",
        "let me get", "let me search", "i will search"
    )
    return any(marker in lowered for marker in planning_markers)


# ══════════════════════════════════════════════════════════════════════════
#  RICH CONSOLE  (panels, tables, non-streaming output)
# ══════════════════════════════════════════════════════════════════════════

THEME = Theme({
    "agent":     "bold #7DCFFF",
    "user":      "bold #9ECE6A",
    "tool.name": "#BB9AF7",
    "tool.ok":   "#9ECE6A",
    "tool.err":  "#F7768E",
    "tool.info": "#73DACA",
    "dim":       "#565F89",
    "header":    "bold #7AA2F7",
    "accent":    "#FF9E64",
    "border":    "#3B4261",
})
console = Console(theme=THEME, highlight=False)


# ══════════════════════════════════════════════════════════════════════════
#  CONFIG & PATHS
# ══════════════════════════════════════════════════════════════════════════

CONFIG_DIR = Path.home() / ".config" / "open-agent"
SESSIONS_DIR = CONFIG_DIR / "sessions"
CONFIG_FILE = CONFIG_DIR / "config.json"
SOUL_FILE = Path(__file__).parent.parent / "SOUL.md"
MEMORY_FILE = CONFIG_DIR / "MEMORY.md"
USER_FILE = CONFIG_DIR / "USER.md"

CONFIG_DIR.mkdir(parents=True, exist_ok=True)
SESSIONS_DIR.mkdir(parents=True, exist_ok=True)

DEFAULT_CONFIG = {
    "base_url": "http://localhost:8083/v1",
    "model_name": "llama.cpp",
    "searxng_url": "http://localhost:8081/search",
    "obsidian_path": "",
    "max_pairs": 3,
    "ctx_limit": 32_768,    # Default for most 8B models; auto-detected from Ollama API if available
    "ctx_warn": 24_000,
    "ctx_compress": 28_000,
    "setup_done": False,
    "memory_nudge_interval": 10,
    "max_summary_tokens":    600,       # was 800 — richer summaries
    "max_tool_result_chars": 6_000,     # was 12000 — tighter per-tool cap
    "use_parallel_loop": True,  # DEFAULT: use raw-OAI loop with parallel tool calls
}

def load_config() -> dict:
    if CONFIG_FILE.exists():
        try:
            return {**DEFAULT_CONFIG, **json.loads(CONFIG_FILE.read_text())}
        except Exception:
            pass
    return DEFAULT_CONFIG.copy()

def save_config(cfg: dict):
    CONFIG_FILE.write_text(json.dumps(cfg, indent=2, ensure_ascii=False))

CFG = load_config()

SOUL_TRIGGERS = {
    "research", "plan", "how should i", "step by step", "steps to",
    "compare", "analyze", "analyse", "strategy", "deep dive",
    "explain in detail", "walk me through", "pros and cons",
    "investigate", "outline", "draft", "write a", "help me build",
}
GROUNDING_TRIGGERS = {
    "today", "now", "current", "latest", "recent", "who is", "what is",
    "price", "weather", "news", "2024", "2025", "this year",
    "this week", "yesterday", "released", "announced",
}

# ── Path safety ───────────────────────────────────────────────────────────────
_PATH_SEP = "/"  # portable path separator

def _validate_path(path: str) -> None:
    """
    Reject paths containing null bytes, newlines, or traversal attempts.
    These are never valid file paths and are a common injection vector.
    """
    if "\x00" in path:
        raise PermissionError("Null byte in path — rejected for safety.")
    if "\n" in path or "\r" in path:
        raise PermissionError("Newline in path — rejected for safety.")
    # Block explicit traversal attempts
    clean = path.replace("\\", "/")
    if ".." in clean:
        raise PermissionError(f"Path traversal '..' in '{path}' — rejected for safety.")


def _safe_expand_path(path: str) -> Path:
    """Resolve path, restrict to user's home, and reject symlink escapes."""
    _validate_path(path)
    p = Path(path).expanduser().resolve()
    home = Path.home().resolve()
    config = CONFIG_DIR.resolve()

    # Follow symlinks and re-check resolved path
    if p.is_symlink():
        p = p.resolve()
        resolved = str(p)
        if not (resolved.startswith(str(home)) or resolved.startswith(str(config))):
            raise PermissionError(
                f"Symlink target '{p}' escapes home directory — rejected for safety."
            )

    if not (str(p).startswith(str(home)) or str(p).startswith(str(config))):
        raise PermissionError(
            f"⛔ Path '{path}' is outside your home directory for safety.\n"
            f"File tools are restricted to ~/... only."
        )
    return p

# ── Terminal safety: block dangerous shell patterns ───────────────────────
_BLOCKED_PATTERNS = [
    # System destruction
    r"\brm\s+-rf\s+/",
    r"\brm\s+--no-preserve-root",
    r"\bmkfs\b",
    r"\bdd\b.*of=/dev/(sd|nvme|hd|vd)",
    r">\s*/dev/sd",
    r":\(\)\s*\{.*\}\s*;",
    # Privilege escalation
    r"\bsudo\s+su\b",
    r"\bsudo\s+bash\b",
    r"\bsudo\s+sh\b",
    r"\bchmod\s+777\s+/",
    r"\bchown\b.*\s+/etc/",
    # Exfiltration / reverse shells
    r"(curl|wget)\s+.*\|\s*(ba)?sh",
    r"\bnc\b.*-e\s+/bin/(ba)?sh",
    r"\bpython\b.*-c.*socket.*connect",
    r"exec\s+\d+<>/dev/tcp",
    r"\bbash\s+-i\s+>&\s+/dev/tcp",
    # Credential theft
    r"cat\s+/etc/(passwd|shadow|sudoers)",
    r"/etc/ssh/.*id_rsa",
    # Home directory nuking
    r"\brm\s+-r[rf]?\s+~",
    r"\brm\s+-r[rf]?\s+\$HOME",
    r"\brm\s+-r[rf]?\s+/home",
    r"\bshred\b",
    r"\bwipefs\b",
    # Fork bombs
    r"\n:(){ :|:& };:",   # newline-variant fork bomb
    r":\(\){:\|:&};:",   # compact fork bomb
    # Systemd / service abuse
    r"\bsystemctl\s+stop\s+(NetworkManager|sshd|cron)",
    r"\bservice\s+(iptables|firewalld|ufw)\s+stop",
    # Package manager override
    r"\bpip\s+install\b.*--break-system",
    r"\bapt-get\s+install\b.*\(.*\)",  # install to custom dir
    # SSH keygen / known_hosts manipulation
    r"\bssh-keygen\b.*-f\s+/",
    # Kernel module tampering
    r"\bmodprobe\s+-r\b",
    r"\binsmod\b",
    # Container escape
    r"docker\s+run\s+--privileged",
    r"nsenter\b",
]

_BLOCKED_RE = [re.compile(p, re.IGNORECASE) for p in _BLOCKED_PATTERNS]

# ── Sanitized subprocess environment ────────────────────────────────────────
# Keys stripped from subprocess env to prevent credential exfiltration
_SANITIZE_ENV_KEYS = frozenset({
    "AWS_SECRET_ACCESS_KEY", "AWS_SESSION_TOKEN",
    "GITHUB_TOKEN", "GITHUB_API_KEY",
    "OPENAI_API_KEY", "ANTHROPIC_API_KEY",
    "HF_TOKEN", "HUGGING_FACE_HUB_TOKEN",
    "DATABASE_URL", "DB_PASSWORD",
    "SECRET_KEY", "SECRET_TOKEN",
    "STRIPE_SECRET_KEY", "PAYMENT_KEY",
    "PRIVATE_KEY", "SSH_PRIVATE_KEY",
})

def _safe_env() -> dict:
    """Return a sanitized environment — strip known secret keys, keep PATH/LANG."""
    env = {}
    for k, v in os.environ.items():
        if k in _SANITIZE_ENV_KEYS:
            env[k] = "[REDACTED]"
        else:
            env[k] = v
    # Always set safety flags
    env.setdefault("PYTHONDONTWRITEBYTECODE", "1")
    env.setdefault("PYTHONUNBUFFERED", "1")
    return env

def _is_dangerous(cmd: str) -> tuple[bool, str]:
    for pattern in _BLOCKED_RE:
        if pattern.search(cmd):
            return True, pattern.pattern
    return False, ""


# ── Tool argument validators ───────────────────────────────────────────────

def _validate_path_arg(args: dict) -> None:
    """Validate that a 'path' arg doesn't contain injection or traversal."""
    path = args.get("path", "")
    if not path:
        raise ValueError("path argument is required")
    if "\x00" in path or "\n" in path or "\r" in path:
        raise ValueError("Invalid characters in path — rejected")
    if ".." in path.replace("\\", "/"):
        raise ValueError("Path traversal '..' — rejected")


def _validate_command_arg(args: dict) -> None:
    """Validate command doesn't match any blocked dangerous pattern."""
    cmd = args.get("command", "")
    if not cmd:
        raise ValueError("command argument is required")
    dangerous, _ = _is_dangerous(cmd)
    if dangerous:
        raise ValueError("Command blocked for safety")


# ── Tool Registry bootstrap ─────────────────────────────────────────────────

def _build_tool_registry() -> ToolRegistry:
    """Build and populate the global tool registry with all known tools."""
    registry = ToolRegistry()

    for _name, _desc, _danger, _examples, _cats, _val in (
        ("web_search",           "search web, returns JSON with title/url/snippet", 0,
         ("latest news", "weather in city"),          ("web",),         None),
        ("cached_web_search",   "same as web_search but cached 1hr",  0,
         ("revisit topic",),                          ("web",),         None),
        ("smart_research",      "parallel web search (up to 4 queries)", 0,
         ("multi-angle research",),                    ("web",),         None),
        ("fetch_page",          "extract readable text from URL",       0,
         ("get article content",),                    ("web",),         None),
        ("run_terminal",        "run shell command, 30s timeout",     1,
         ("git status", "ls -la"),                  ("terminal",),     _validate_command_arg),
        ("read_file",           "read local file (max 8,000 chars)",  0,
         ("cat ~/.bashrc", "view config"),          ("file",),         _validate_path_arg),
        ("write_file",         "write/create local file",             1,
         ("write to ~/notes.txt",),                 ("file",),         _validate_path_arg),
        ("load_soul",          "load SOUL.md behavioral instructions", 0,
         ("load soul",),                            ("soul",),         None),
        ("search_obsidian",     "search Obsidian vault notes",         0,
         ("search notes for X",),                  ("obsidian",),      None),
        ("read_obsidian_note", "read specific Obsidian note",         0,
         ("read note X",),                         ("obsidian",),     None),
        ("write_obsidian_note","create/append Obsidian note",        1,
         ("save to obsidian",),                    ("obsidian",),     None),
        ("update_memory",      "append to MEMORY.md for cross-session recall", 1,
         ("remember that X",),                    ("memory",),        None),
        ("read_memory",        "read MEMORY.md",                     0,
         ("what do you remember",),               ("memory",),        None),
        ("update_user_profile","update USER.md profile",              1,
         ("update my profile",),                  ("memory",),        None),
        ("search_sessions",    "search past conversations via FTS5", 0,
         ("what did we discuss about X",),         ("memory",),        None),
        ("create_pptx",        "create PowerPoint from slide list",  1,
         ("make slides about X",),                 ("office",),        None),
        ("read_rss_by_name",   "fetch RSS feed items",              0,
         ("show hacker news",),                    ("web",),           None),
    ):
        spec = ToolSpec(
            name=_name, description=_desc, danger=_danger,
            examples=_examples, categories=_cats, validator=_val,
        )
        registry.register(spec)

    return registry


# Populate the global registry at import time
TOOL_REGISTRY = _build_tool_registry()


# ══════════════════════════════════════════════════════════════════════════
#  AGENT SETUP
# ══════════════════════════════════════════════════════════════════════════

provider = OpenAIProvider(base_url=CFG["base_url"], api_key="lm-studio")
model    = OpenAIChatModel(model_name=CFG["model_name"], provider=provider)

SYSTEM_PROMPT = """You are a fast, grounded local AI agent with chain-of-thought reasoning.

CHAIN-OF-THINK PROCESS (think before acting):
Before calling ANY tool, show your thinking using <reasoning></reasoning> tags.
Example:
  <think>I need to check the current date to answer this time-sensitive question. I'll use run_terminal with the date command.</think>
  Then make the tool call.

For complex multi-step tasks, break down your reasoning:
1. WHAT do I need? (goal)
2. WHICH tool gives me that? (selection)
3. HOW will I use the result? (execution)
4. WHAT if it fails? (error handling)

RITUALS — before any real-world or time-sensitive query:
1. run_terminal("date && uname -r") — anchor current time/system
2. web_search if knowledge might be stale (default: assume it is for recent topics)
3. Confirm facts across >=2 results. ALWAYS include direct URLs in responses.
TOOL DESCRIPTIONS (use as reference):

web_search(query) → search web, returns JSON with title/url/content. For current facts, news, prices.

cached_web_search(query) → same but caches 1hr — use for topics you might revisit.

smart_research(queries) → parallel web search (up to 4), merges results. For multi-angle research.

fetch_page(url) → extract readable text from URL. Use after web_search finds promising links.

run_terminal(command) → run shell command. TIMEOUT: 30s. DANGEROUS PATTERNS ARE BLOCKED.
  • On error: the error will include a HINT on how to fix it
  • On NameError: define variable first, retry
  • On ImportError: install missing package, retry
  • On FileNotFoundError: check path or create parent dirs
  • On PermissionError: use different path in ~/

read_file(path) → read local file (up to 8,000 chars). RESTRICTED TO ~/ FOR SAFETY.

write_file(path, content) → write local file. RESTRICTED TO ~/ FOR SAFETY. Auto-creates parent dirs.

create_pptx(path, slides) → create PowerPoint. slides=[{"title":"X","content":"Y"},...]

update_memory(content) → append to persistent MEMORY.md for cross-session recall

read_memory() → read MEMORY.md

search_sessions(query) → search past conversations via FTS5

search_obsidian(query) → search Obsidian vault

read_obsidian_note(note_name) → read specific Obsidian note

write_obsidian_note(note_name, content) → create/append Obsidian note

load_soul() → load extended SOUL.md behavioral instructions for complex tasks

read_rss_by_name(name) → fetch RSS feeds. names: hn, verge, arxiv_ai, security, etc. categories: tech, ai, engineering

ERROR RECOVERY:
- On NameError: Define the variable first, then retry
- On SyntaxError: Check parentheses/colons, fix and retry

When writing code: use write_file. Show key code in response.
Be direct. Use Markdown. No filler."""

agent = Agent(
    model=model,
    system_prompt=SYSTEM_PROMPT + """

IMPORTANT: Always include raw URLs in responses. Never summarize - list links directly.""",
    deps_type=None,
    model_settings={
        "temperature": 0.8,
        "max_tokens": 8192,
        "ctx_limit": 28000,
        "ctx_warn": 20000,
        "ctx_compress": 24000,
        "max_pairs": 8,
        "max_tool_steps": 15,
        "max_tool_result_chars": 6000,
        "frequency_penalty": 0.35,
        "presence_penalty": 0.15,
    },
    retries=2,
)


# ══════════════════════════════════════════════════════════════════════════
#  TOOL BADGE  (in-place ANSI overwrite — no Rich.Live)
# ══════════════════════════════════════════════════════════════════════════

class ToolBadge:
    def start(self, name: str, detail: str):
        sys.stderr.write(f"  {C_AMBER}⟳{RESET} {C_PURPLE}{name}{RESET} {DIM}{detail[:120]}{RESET}\n")
        sys.stderr.flush()

    def done(self, name: str, summary: str, ok: bool = True, chars: int = 0):
        icon = f"{C_GREEN}✓{RESET}" if ok else f"{C_RED}✗{RESET}"
        ch = f"  {DIM}{chars:,}c{RESET}" if chars else ""
        sys.stderr.write(f"  {icon} {C_PURPLE}{name}{RESET} {DIM}{summary[:120]}{ch}{RESET}\n")
        sys.stderr.flush()

    def info(self, label: str):
        sys.stderr.write(f"  {DIM}· {label}{RESET}\n")
        sys.stderr.flush()

td = ToolBadge()


# ══════════════════════════════════════════════════════════════════════════
#  TOOLS
# ══════════════════════════════════════════════════════════════════════════

@agent.tool
async def web_search(ctx: RunContext, query: str) -> str:
    """Search the web via SearxNG. Use for current facts, news, or anything potentially stale."""
    td.start("web_search", query)
    try:
        async with httpx.AsyncClient(timeout=12) as c:
            r = await c.get(CFG["searxng_url"],
                params={"q": query, "format": "json", "categories": "general", "language": "en"})
            r.raise_for_status()
        results = r.json().get("results", [])[:7]
        if not results:
            td.done("web_search", "no results", ok=False); return "No results."
        out = [f"[{i}] {x.get('title','')}\nURL: {x.get('url','')}\n{x.get('content','')[:400]}"
               for i, x in enumerate(results, 1)]
        body = "\n\n".join(out)
        td.done("web_search", f"{len(results)} results", chars=len(body))
        return body
    except Exception as e:
        td.done("web_search", str(e), ok=False); return f"Search error: {e}"


@agent.tool
async def cached_web_search(ctx: RunContext, query: str) -> str:
    """Search web with local caching (1hr TTL). Use for queries that may repeat.
    
    Example: cached_web_search(query="ollama models") returns cached if queried before.
    """
    from open_agent.tool_cache import get_cached_result, cache_tool_result
    cached = get_cached_result("web_search", query)
    if cached:
        return f"[CACHED]\n{cached}"
    td.start("cached_web_search", query[:50])
    try:
        async with httpx.AsyncClient(timeout=12) as c:
            r = await c.get(CFG["searxng_url"],
                params={"q": query, "format": "json", "categories": "general"})
            r.raise_for_status()
        results = r.json().get("results", [])[:7]
        if not results:
            return "No results."
        out = [f"[{i}] {x.get('title','')}\n{x.get('url','')}\n{x.get('content','')[:300]}"
               for i, x in enumerate(results, 1)]
        body = "\n\n".join(out)
        cache_tool_result("web_search", query, body)
        return body
    except Exception as e:
        return f"Search error: {e}"


@agent.tool
async def smart_research(ctx: RunContext, queries: list[str]) -> str:
    """Run up to 4 parallel web queries and merge. For multi-angle research."""
    queries = queries[:4]
    td.start("smart_research", " | ".join(queries))
    async def _one(q: str) -> str:
        try:
            async with httpx.AsyncClient(timeout=12) as c:
                r = await c.get(CFG["searxng_url"],
                    params={"q": q, "format": "json", "categories": "general", "language": "en"})
                r.raise_for_status()
            res = r.json().get("results", [])[:4]
            return "\n".join([f"### {q}"] +
                [f"**{i.get('title','')}**\n{i.get('url','')}\n{i.get('content','')[:350]}" for i in res]
            ) if res else f"### {q}\nNo results."
        except Exception as e: return f"### {q}\nError: {e}"
    parts  = await asyncio.gather(*[_one(q) for q in queries])
    merged = "\n\n---\n\n".join(parts)
    td.done("smart_research", f"merged {len(queries)} queries", chars=len(merged))
    return merged


@agent.tool
async def fetch_page(ctx: RunContext, url: str) -> str:
    """Fetch readable text from a URL."""
    td.start("fetch_page", url[:70])
    try:
        async with httpx.AsyncClient(timeout=15, follow_redirects=True) as c:
            r = await c.get(url, headers={"User-Agent": "Mozilla/5.0"})
            r.raise_for_status()
        text = re.sub(r"<[^>]+>", " ", r.text)
        text = re.sub(r"\s{3,}", "\n\n", text).strip()[:6000]
        td.done("fetch_page", url[:50], chars=len(text)); return text
    except Exception as e:
        td.done("fetch_page", str(e), ok=False); return f"Fetch error: {e}"


@agent.tool
async def run_terminal(ctx: RunContext, command: str) -> str:
    """Run a safe shell command. Dangerous patterns are blocked for your safety.
    
    On error: The model will see the error with a HINT on how to fix it.
    - NameError → define variable first, retry
    - ImportError → install missing package, retry  
    - FileNotFoundError → check path or create parent dirs
    - PermissionError → use different path in ~/
    - TimeoutExpired → optimize command or increase timeout
    """
    dangerous, pattern = _is_dangerous(command)
    if dangerous:
        td.done("run_terminal", f"BLOCKED — dangerous pattern", ok=False)
        return (
            f"⛔ Command blocked for safety.\n"
            f"Pattern matched: {pattern}\n"
            f"If this is a legitimate command, run it manually in your terminal."
        )
    td.start("run_terminal", command[:70])

    # Auto-retry on retryable errors (up to 2 retries)
    for attempt in range(3):
        try:
            proc = subprocess.run(
                command, shell=True, capture_output=True, text=True, timeout=30,
                env=_safe_env(),
            )
            ok = proc.returncode == 0
            out = proc.stdout.strip(); err = proc.stderr.strip()

            # If error occurred, add fix hint
            if not ok:
                error_output = err or f"Exit code {proc.returncode}"
                analysis = _analyze_error(error_output)
                if analysis:
                    td.done("run_terminal", f"error + hint", ok=False)
                    return (
                        f"STDOUT:\n{proc.stdout}\n"
                        f"STDERR:\n{error_output}\n"
                        f"Exit: {proc.returncode}\n\n"
                        f"💡 {analysis}"
                    )
                td.done("run_terminal", (out or err or "done")[:60], ok=ok)
                return f"STDOUT:\n{proc.stdout}\nSTDERR:\n{proc.stderr}\nExit: {proc.returncode}"
            else:
                td.done("run_terminal", (out or err or "done")[:60], ok=ok)
                return f"STDOUT:\n{proc.stdout}\nSTDERR:\n{proc.stderr}\nExit: {proc.returncode}"

        except subprocess.TimeoutExpired:
            if attempt < 2:
                td.info(f"timeout — retrying ({attempt+1}/3)")
                await asyncio.sleep(1)
                continue
            td.done("run_terminal", "timed out (30s)", ok=False)
            return "TIMEOUT: Command exceeded 30s. Optimize the command or use a simpler approach."
        except Exception as e:
            hint = _analyze_error(str(e))
            td.done("run_terminal", str(e)[:60], ok=False)
            if hint:
                return f"ERROR: {e}\n\n💡 {hint}"
            return f"Error: {e}"

    return "Failed after 3 attempts."


@agent.tool
async def read_file(ctx: RunContext, path: str) -> str:
    """Read a local file (up to 8,000 chars). RESTRICTED TO ~/ FOR SAFETY.

    On error: Check path, file exists, or use a different path.
    """
    td.start("read_file", path)
    for attempt in range(3):
        try:
            p = _safe_expand_path(path)
            content = p.read_text(encoding="utf-8")[:8000]
            td.done("read_file", str(p), chars=len(content))
            return content
        except PermissionError as e:
            td.done("read_file", str(e), ok=False)
            return f"Permission denied: {e}\n\n💡 Hint: Use a path you have read access to, or check file permissions."
        except FileNotFoundError as e:
            td.done("read_file", "not found", ok=False)
            return f"File not found: {path}\n\n💡 Hint: Check the path is correct. File paths are case-sensitive on Linux."
        except Exception as e:
            if attempt < 2:
                td.info(f"read error — retrying")
                await asyncio.sleep(0.1)
                continue
            td.done("read_file", str(e), ok=False)
            return f"Read error: {e}\n\n💡 Hint: Check if the file is readable or try a different path."

@agent.tool
async def write_file(ctx: RunContext, path: str, content: str, expand_user: bool = True) -> str:
    """Write content to a local file. RESTRICTED TO ~/ FOR SAFETY.
    
    On error: Model should retry with correct path or different location.
    - PermissionError → use different path in ~/
    - FileNotFoundError → create parent dirs or use existing path
    """
    td.start("write_file", f"{path} ({len(content):,} chars)")

    # Auto-retry with mkdir on permission/path errors
    for attempt in range(3):
        try:
            p = _safe_expand_path(path)
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(content, encoding="utf-8")
            td.done("write_file", f"→ {p}", chars=len(content))
            return f"✓ Written {len(content):,} chars to {p}"
        except PermissionError as e:
            td.done("write_file", str(e), ok=False)
            return f"Permission denied: {e}\n\n💡 Hint: Use a different path in your home directory."
        except FileNotFoundError as e:
            if attempt < 2:
                td.info(f"path error — retrying with mkdir")
                try:
                    p2 = _safe_expand_path(path)
                    p2.parent.mkdir(parents=True, exist_ok=True)
                    p2.write_text(content, encoding="utf-8")
                    td.done("write_file", f"→ {p2} (auto-mkdir)", chars=len(content))
                    return f"✓ Written {len(content):,} chars to {p2}"
                except Exception:
                    pass
            td.done("write_file", str(e), ok=False)
            return f"Path error: {e}\n\n💡 Hint: Create parent directories first or use a path that exists."
        except Exception as e:
            td.done("write_file", str(e), ok=False)
            hint = _analyze_error(str(e))
            if hint:
                return f"Write error: {e}\n\n💡 {hint}"
            return f"Write error: {e}"

    return "Failed after 3 attempts."

@agent.tool
async def load_soul(ctx: RunContext) -> str:
    """Load extended behavioral instructions from SOUL.md. Call for complex/research tasks."""
    td.start("load_soul", "SOUL.md")
    try:
        content = SOUL_FILE.read_text(encoding="utf-8")
        td.done("load_soul", "behavioral core loaded", chars=len(content)); return content
    except FileNotFoundError:
        td.done("load_soul", "SOUL.md not found", ok=False)
        return "SOUL.md not found alongside agent.py."
    except Exception as e:
        td.done("load_soul", str(e), ok=False); return f"Error: {e}"


# ── Obsidian ──────────────────────────────────────────────────────────────

def _vault() -> Optional[Path]:
    p = CFG.get("obsidian_path", "")
    return Path(p) if p else None


@agent.tool
async def search_obsidian(ctx: RunContext, query: str) -> str:
    """Search your Obsidian vault for notes containing specific text."""
    vault = _vault()
    if not vault or not vault.exists():
        return "Obsidian vault not configured. Run /setup to configure."
    td.start("search_obsidian", query)
    try:
        results = []
        for p in vault.rglob("*.md"):
            if any(part.startswith('.') for part in p.parts): continue
            if query.lower() in p.read_text(encoding="utf-8", errors="ignore").lower():
                results.append(f"- {p.relative_to(vault)}")
        if not results:
            td.done("search_obsidian", "no matches", ok=False); return "No matching notes."
        td.done("search_obsidian", f"{len(results)} matches")
        return "Matching notes:\n" + "\n".join(results[:25])
    except Exception as e:
        td.done("search_obsidian", str(e), ok=False); return f"Error: {e}"


@agent.tool
async def read_obsidian_note(ctx: RunContext, note_name: str) -> str:
    """Read a specific note from your Obsidian vault."""
    vault = _vault()
    if not vault or not vault.exists():
        return "Obsidian vault not configured. Run /setup."
    if not note_name.endswith(".md"): note_name += ".md"
    td.start("read_obsidian_note", note_name)
    try:
        matches = [m for m in vault.rglob(note_name)
                   if not any(p.startswith('.') for p in m.parts)]
        if not matches:
            td.done("read_obsidian_note", "not found", ok=False)
            return f"Note '{note_name}' not found."
        content = matches[0].read_text(encoding="utf-8")[:10000]
        td.done("read_obsidian_note", note_name, chars=len(content)); return content
    except Exception as e:
        td.done("read_obsidian_note", str(e), ok=False); return f"Error: {e}"


@agent.tool
async def write_obsidian_note(ctx: RunContext, note_name: str, content: str, append: bool = True) -> str:
    """Create or append to an Obsidian note."""
    vault = _vault()
    if not vault or not vault.exists():
        return "Obsidian vault not configured. Run /setup."
    if not note_name.endswith(".md"): note_name += ".md"
    td.start("write_obsidian_note", f"{note_name} append={append}")
    try:
        matches = [m for m in vault.rglob(note_name)
                   if not any(p.startswith('.') for p in m.parts)]
        target = matches[0] if matches else vault / note_name
        mode   = "a" if append and target.exists() else "w"
        with open(target, mode, encoding="utf-8") as f:
            f.write(("\n\n" if mode == "a" else "") + content)
        action = "Appended to" if mode == "a" else "Created"
        td.done("write_obsidian_note", f"{action} {note_name}")
        return f"✓ {action} {target.relative_to(vault)}"
    except Exception as e:
        td.done("write_obsidian_note", str(e), ok=False); return f"Error: {e}"


# ── Persistent Memory Tools (NEW) ──────────────────────────────────────────

@agent.tool
async def update_memory(ctx: RunContext, content: str, append: bool = True) -> str:
    """Update persistent MEMORY.md. Appends by default, use append=false to replace."""
    td.start("update_memory", f"{len(content)} chars")
    try:
        MEMORY_FILE.parent.mkdir(parents=True, exist_ok=True)
        if append and MEMORY_FILE.exists():
            existing = MEMORY_FILE.read_text(encoding="utf-8")
            MEMORY_FILE.write_text(existing + "\n\n" + content, encoding="utf-8")
        else:
            MEMORY_FILE.write_text(content, encoding="utf-8")
        td.done("update_memory", f"updated · append={append}")
        return f"✓ MEMORY.md {'appended' if append else 'written'}"
    except Exception as e:
        td.done("update_memory", str(e), ok=False)
        return f"Error: {e}"


@agent.tool
async def read_memory(ctx: RunContext) -> str:
    """Read persistent MEMORY.md."""
    td.start("read_memory", "MEMORY.md")
    try:
        if not MEMORY_FILE.exists():
            td.done("read_memory", "not found", ok=False)
            return "MEMORY.md not found. Use update_memory to create it."
        content = MEMORY_FILE.read_text(encoding="utf-8", errors="ignore")[:3000]
        td.done("read_memory", f"{len(content)} chars")
        return content
    except Exception as e:
        td.done("read_memory", str(e), ok=False)
        return f"Error: {e}"


@agent.tool
async def update_user_profile(ctx: RunContext, content: str, append: bool = False) -> str:
    """Update USER.md with user info (name, preferences, context)."""
    td.start("update_user_profile", f"{len(content)} chars")
    try:
        USER_FILE.parent.mkdir(parents=True, exist_ok=True)
        if append and USER_FILE.exists():
            existing = USER_FILE.read_text(encoding="utf-8")
            USER_FILE.write_text(existing + "\n\n" + content, encoding="utf-8")
        else:
            USER_FILE.write_text(content, encoding="utf-8")
        td.done("update_user_profile", "updated")
        return "✓ USER.md updated"
    except Exception as e:
        td.done("update_user_profile", str(e), ok=False)
        return f"Error: {e}"


@agent.tool
async def search_sessions(ctx: RunContext, query: str, limit: int = 10) -> str:
    """Search all past sessions via FTS5. Returns matching conversation pairs."""
    td.start("search_sessions", query[:30])
    try:
        results = _db.search(query, limit=limit)
        if not results:
            td.done("search_sessions", "no matches", ok=False)
            return f"No matches for: {query}"
        out = [f"**Search: {query}** ({len(results)} matches)\n"]
        for r in results[:limit]:
            sid = r.get("session_id", "?")[:12]
            title = r.get("title", "")[:40]
            user = r.get("user", r.get("user_msg", ""))[:80].replace("\n", " ")
            assist = r.get("assist", r.get("assist_msg", ""))[:120].replace("\n", " ")
            out.append(f"---")
            out.append(f"Session: {sid} · {title}")
            out.append(f"User: {user}")
            out.append(f"Assistant: {assist}")
        result = "\n".join(out)
        td.done("search_sessions", f"{len(results)} matches", chars=len(result))
        return result
    except Exception as e:
        td.done("search_sessions", str(e), ok=False)
        return f"Error: {e}"


@agent.tool
async def create_pptx(ctx: RunContext, path: str, slides: list[dict]) -> str:
    """Create PowerPoint presentation from slide data.
    
    Args:
        path: Output .pptx file path
        slides: List of slide dicts with 'title' and optional 'content', 'bullets'
    
    Example slides:
        [{"title": "Intro", "content": "Welcome"}, {"title": "Agenda", "bullets": ["Point 1", "Point 2"]}]
    """
    td.start("create_pptx", path)
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(6)
        
        for slide_data in slides:
            sp = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
            title = sp.shapes.title
            title.text = slide_data.get("title", "")
            
            # Content or bullets
            content_body = sp.placeholders[1]
            tf = content_body.text_frame
            tf.clear()
            
            if "content" in slide_data:
                tf.text = slide_data["content"]
            elif "bullets" in slide_data:
                for i, bullet in enumerate(slide_data["bullets"]):
                    if i == 0:
                        p = tf.paragraphs[0]
                        p.text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0
        
        prs.save(path)
        td.done("create_pptx", f"created {len(slides)} slides")
        return f"✓ Created {path} with {len(slides)} slides"
    except ImportError:
        td.done("create_pptx", "python-pptx not installed", ok=False)
        return "Error: python-pptx not installed. Run: pip install python-pptx"
    except Exception as e:
        td.done("create_pptx", str(e), ok=False)
        return f"Error: {e}"


# ── RSS ───────────────────────────────────────────────────────────────────

RSS_SOURCES: dict[str, str] = {
    "eff": "https://www.eff.org/rss/updates.xml",
    "schneier": "https://www.schneier.com/feed/atom/",
    "krebs": "https://krebsonsecurity.com/feed/",
    "bleeping": "https://www.bleepingcomputer.com/feed/",
    "techcrunch": "https://techcrunch.com/feed/",
    "verge": "http://www.theverge.com/rss/full.xml",
    "engadget": "http://www.engadget.com/rss-full.xml",
    "venturebeat": "http://venturebeat.com/feed/",
    "ars": "https://arstechnica.com/feed/",
    "hn": "https://news.ycombinator.com/rss",
    "hn_show": "http://hnrss.org/show",
    "hn_launches": "https://hnrss.org/launches",
    "pragmatic": "https://blog.pragmaticengineer.com/rss/",
    "cloudflare": "https://blog.cloudflare.com/rss/",
    "stripe": "https://stripe.com/blog/feed.rss",
    "meta_eng": "https://engineering.fb.com/feed/",
    "julia_evans": "https://jvns.ca/atom.xml",
    "danluu": "https://danluu.com/atom.xml",
    "arxiv_ai": "https://rss.arxiv.org/rss/cs.AI",
    "arxiv_cl": "https://rss.arxiv.org/rss/cs.CL",
    "huggingface": "https://huggingface.co/blog/feed.xml",
    "farnam": "https://fs.blog/feed/",
    "producthunt": "http://www.producthunt.com/feed",
}
BROAD_CATS: dict[str, list[str]] = {
    "tech":        ["techcrunch", "verge", "ars", "venturebeat", "hn"],
    "technology":  ["techcrunch", "verge", "ars", "venturebeat", "hn"],
    "news":        ["techcrunch", "verge", "ars", "hn"],
    "startups":    ["techcrunch", "hn_launches", "producthunt"],
    "engineering": ["pragmatic", "cloudflare", "stripe", "meta_eng"],
    "ai":          ["arxiv_ai", "arxiv_cl", "huggingface"],
    "security":    ["eff", "schneier", "krebs", "bleeping"],
}

async def _fetch_rss(url: str, limit: int = 8) -> str:
    feed = feedparser.parse(url)
    if feed.bozo: return f"RSS error: {feed.bozo_exception}"
    if not feed.entries: return f"No entries: {url}"
    out = [f"**{feed.feed.get('title','Feed')}**\n"]
    for i, e in enumerate(feed.entries[:min(limit, 15)], 1):
        out.append(f"{i}. **{e.get('title','No title')}**\n"
                   f"   {e.get('published', e.get('updated',''))}\n"
                   f"   {e.get('link','')}\n"
                   f"   {e.get('summary', e.get('description',''))[:260].strip()}\n")
    return "\n".join(out)

@agent.tool
async def read_rss_by_name(ctx: RunContext, name: str, limit: int = 8) -> str:
    """Fetch RSS from curated sources. Names: hn, verge, arxiv_ai... Categories: tech, ai, security, engineering, startups."""
    td.start("read_rss_by_name", name)
    n = name.lower().strip()
    if n in BROAD_CATS:
        results = []
        for fn in BROAD_CATS[n]:
            if fn in RSS_SOURCES:
                try: results.append(f"### {fn.upper()}\n{await _fetch_rss(RSS_SOURCES[fn], 4)}")
                except Exception: pass
        merged = "\n\n---\n\n".join(results)
        td.done("read_rss_by_name", f"merged {len(results)} feeds", chars=len(merged))
        return merged or f"No content for '{name}'."
    url = RSS_SOURCES.get(n) or next((v for k, v in RSS_SOURCES.items() if n in k or k in n), None)
    if url:
        c = await _fetch_rss(url, limit)
        td.done("read_rss_by_name", name, chars=len(c)); return c
    td.done("read_rss_by_name", f"unknown: '{name}'", ok=False)
    return f"Unknown: '{name}'. Categories: {', '.join(BROAD_CATS)}. Sources: {', '.join(sorted(RSS_SOURCES))[:200]}..."


# ══════════════════════════════════════════════════════════════════════════
#  SQLITE SESSION DATABASE  (NEW — enhances JSON storage)
# ══════════════════════════════════════════════════════════════════════════

class SessionDB:
    """
    SQLite session storage with FTS5 full-text search.
    Maintains JSON files for backward compat — new data goes to SQLite.
    """
    _instance = None
    _lock = threading.Lock()

    def __init__(self, db_path: Path):
        self.path = db_path
        self.conn: Optional[sqlite3.Connection] = None
        self._init_db()

    @classmethod
    def get(cls, db_path: Optional[Path] = None) -> "SessionDB":
        """Singleton to avoid duplicate connections."""
        if cls._instance is None:
            with cls._lock:
                if cls._instance is None:
                    cls._instance = cls(db_path or (CONFIG_DIR / "sessions.db"))
        return cls._instance

    def _init_db(self):
        """Initialize database with tables and FTS5."""
        db = self._get_conn()
        db.execute("PRAGMA journal_mode=WAL")
        db.execute("PRAGMA synchronous=NORMAL")
        db.execute("PRAGMA cache_size=-32000")  # 32MB cache

        db.execute("""
            CREATE TABLE IF NOT EXISTS sessions (
                id TEXT PRIMARY KEY,
                title TEXT,
                created_at TEXT NOT NULL,
                last_active TEXT NOT NULL,
                summary TEXT DEFAULT '',
                memory_md TEXT DEFAULT '',
                user_md TEXT DEFAULT '',
                turn_count INTEGER DEFAULT 0
            )
        """)

        db.execute("""
            CREATE TABLE IF NOT EXISTS pairs (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                session_id TEXT NOT NULL,
                user_msg TEXT NOT NULL,
                assist_msg TEXT NOT NULL,
                created_at TEXT NOT NULL,
                FOREIGN KEY (session_id) REFERENCES sessions(id) ON DELETE CASCADE
            )
        """)

        db.execute("""
            CREATE VIRTUAL TABLE IF NOT EXISTS pairs_fts USING fts5(
                session_id,
                user_msg,
                assist_msg,
                content='pairs',
                content_rowid='id'
            )
        """)

        # Triggers for FTS sync
        db.execute("""
            CREATE TRIGGER IF NOT EXISTS pairs_ai AFTER INSERT ON pairs BEGIN
                INSERT INTO pairs_fts(rowid, session_id, user_msg, assist_msg)
                VALUES (new.id, new.session_id, new.user_msg, new.assist_msg);
            END
        """)
        db.execute("""
            CREATE TRIGGER IF NOT EXISTS pairs_ad AFTER DELETE ON pairs BEGIN
                INSERT INTO pairs_fts(pairs_fts, rowid, session_id, user_msg, assist_msg)
                VALUES ('delete', old.id, old.session_id, old.user_msg, old.assist_msg);
            END
        """)
        db.commit()

    def _get_conn(self) -> sqlite3.Connection:
        if self.conn is None:
            self.conn = sqlite3.connect(str(self.path), check_same_thread=False)
            self.conn.row_factory = sqlite3.Row
        return self.conn

    def create_session(self, sid: str) -> dict:
        now = datetime.now().isoformat()
        db = self._get_conn()
        db.execute(
            "INSERT OR IGNORE INTO sessions (id, title, created_at, last_active) VALUES (?, ?, ?, ?)",
            (sid, "", now, now)
        )
        db.commit()
        return {"id": sid, "title": "", "created_at": now, "last_active": now}

    def get_session(self, sid: str) -> Optional[dict]:
        db = self._get_conn()
        row = db.execute("SELECT * FROM sessions WHERE id = ?", (sid,)).fetchone()
        return dict(row) if row else None

    def list_sessions(self, limit: int = 50) -> list[dict]:
        db = self._get_conn()
        rows = db.execute(
            "SELECT id, title, created_at, last_active, turn_count, summary FROM sessions ORDER BY last_active DESC LIMIT ?",
            (limit,)
        ).fetchall()
        return [dict(r) for r in rows]

    def add_pair(self, sid: str, user_msg: str, assist_msg: str) -> None:
        now = datetime.now().isoformat()
        db = self._get_conn()
        db.execute(
            "INSERT INTO pairs (session_id, user_msg, assist_msg, created_at) VALUES (?, ?, ?, ?)",
            (sid, user_msg, assist_msg, now)
        )
        db.execute(
            "UPDATE sessions SET last_active = ?, turn_count = turn_count + 1 WHERE id = ?",
            (now, sid)
        )
        db.commit()

    def get_pairs(self, sid: str) -> List[Tuple[str, str]]:
        db = self._get_conn()
        rows = db.execute(
            "SELECT user_msg, assist_msg FROM pairs WHERE session_id = ? ORDER BY id",
            (sid,)
        ).fetchall()
        return [(r["user_msg"], r["assist_msg"]) for r in rows]

    def get_pair_count(self, sid: str) -> int:
        db = self._get_conn()
        row = db.execute(
            "SELECT COUNT(*) as cnt FROM pairs WHERE session_id = ?",
            (sid,)
        ).fetchone()
        return row["cnt"] if row else 0

    def update_summary(self, sid: str, summary: str) -> None:
        db = self._get_conn()
        db.execute("UPDATE sessions SET summary = ? WHERE id = ?", (summary, sid))
        db.commit()

    def update_memory(self, sid: str, memory_md: str = None, user_md: str = None) -> None:
        db = self._get_conn()
        if memory_md is not None:
            db.execute("UPDATE sessions SET memory_md = ? WHERE id = ?", (memory_md, sid))
        if user_md is not None:
            db.execute("UPDATE sessions SET user_md = ? WHERE id = ?", (user_md, sid))
        db.commit()

    def get_memory(self, sid: str) -> Tuple[str, str]:
        db = self._get_conn()
        row = db.execute(
            "SELECT memory_md, user_md FROM sessions WHERE id = ?",
            (sid,)
        ).fetchone()
        if row:
            return row["memory_md"] or "", row["user_md"] or ""
        return "", ""

    def search(self, query: str, limit: int = 20) -> List[dict]:
        """Full-text search via FTS5."""
        db = self._get_conn()
        try:
            rows = db.execute("""
                SELECT p.session_id, s.title, p.user_msg, p.assist_msg,
                       snippet(pairs_fts, 1, '[', ']', '...', 32) as user_snippet,
                       snippet(pairs_fts, 2, '[', ']', '...', 32) as assist_snippet
                FROM pairs_fts p
                JOIN sessions s ON p.session_id = s.id
                WHERE pairs_fts MATCH ?
                ORDER BY rank
                LIMIT ?
            """, (query, limit)).fetchall()
            if rows:
                return [dict(r) for r in rows]
        except sqlite3.OperationalError:
            pass

        # Fallback to LIKE search
        like_q = f"%{query}%"
        rows = db.execute("""
            SELECT p.session_id, s.title, p.user_msg, p.assist_msg
            FROM pairs p JOIN sessions s ON p.session_id = s.id
            WHERE p.user_msg LIKE ? OR p.assist_msg LIKE ?
            ORDER BY p.id DESC LIMIT ?
        """, (like_q, like_q, limit)).fetchall()
        return [{"session_id": r["session_id"], "title": r["title"],
                 "user": r["user_msg"][:200], "assist": r["assist_msg"][:200]} for r in rows]

    def delete_session(self, sid: str) -> None:
        db = self._get_conn()
        db.execute("DELETE FROM pairs WHERE session_id = ?", (sid,))
        db.execute("DELETE FROM sessions WHERE id = ?", (sid,))
        db.commit()

    def close(self):
        if self.conn:
            self.conn.close()
            self.conn = None


_db = SessionDB.get()


# ══════════════════════════════════════════════════════════════════════════
#  SESSION MANAGER
# ══════════════════════════════════════════════════════════════════════════

class Session:
    """
    A single conversation session.
    Uses SQLite via SessionDB for storage, with JSON fallback for existing sessions.
    Contains: id, title, created_at, last_active, summary, memory_md, user_md, pairs
    """
    def __init__(self, sid: Optional[str] = None):
        self.id = sid or uuid.uuid4().hex[:12]
        self.path = SESSIONS_DIR / f"{self.id}.json"
        self.pairs: list[tuple[str, str]] = []
        self.summary: str = ""
        self.title: str = ""
        self.memory_md: str = ""
        self.user_md: str = ""
        self.created_at: str = datetime.now().isoformat()
        self.last_active: str = self.created_at
        self._pair_count: int = 0

        # Try SQLite first
        row = _db.get_session(self.id)
        if row:
            self.title = row.get("title", "")
            self.summary = row.get("summary", "")
            self.memory_md = row.get("memory_md", "")
            self.user_md = row.get("user_md", "")
            self.created_at = row.get("created_at", self.created_at)
            self.last_active = row.get("last_active", self.last_active)
            self._pair_count = row.get("turn_count", 0)
            self.pairs = _db.get_pairs(self.id)
        elif sid and self.path.exists():
            # Fallback: load from JSON
            self._load()

    def _load(self):
        """Load from JSON for backward compat."""
        try:
            d = json.loads(self.path.read_text())
            self.pairs = [(p["u"], p["a"]) for p in d.get("pairs", []) if "u" in p]
            self.summary = d.get("summary", "")
            self.title = d.get("title", "")
            self.created_at = d.get("created_at", self.created_at)
            self.last_active = d.get("last_active", self.last_active)
            self._pair_count = len(self.pairs)
        except Exception:
            pass

    def save(self):
        """Save to SQLite and JSON for backward compat."""
        self.last_active = datetime.now().isoformat()
        if not self.title and self.pairs:
            self.title = self.pairs[0][0][:60].replace("\n", " ").strip()
        # Save to SQLite
        if _db.get_session(self.id):
            _db.update_summary(self.id, self.summary)
            _db.update_memory(self.id, self.memory_md, self.user_md)
        else:
            _db.create_session(self.id)
            _db.update_summary(self.id, self.summary)
            _db.update_memory(self.id, self.memory_md, self.user_md)
        # Also save JSON for backward compat
        try:
            self.path.write_text(json.dumps({
                "id": self.id,
                "title": self.title,
                "created_at": self.created_at,
                "last_active": self.last_active,
                "summary": self.summary,
                "pairs": [{"u": u, "a": a} for u, a in self.pairs],
            }, indent=2, ensure_ascii=False))
        except Exception:
            pass

    def delete(self):
        _db.delete_session(self.id)
        self.path.unlink(missing_ok=True)

    def _pair_count_estimate(self) -> int:
        return len(self.pairs)

    def compress(self) -> str:
        """Compress old pairs into summary — v2 with tool metadata preservation.
        
        Token-efficient: keeps tool call patterns for reuse without re-calling.
        """
        max_pairs = CFG["max_pairs"]
        if len(self.pairs) <= max_pairs:
            return "nothing to compress"

        old_pairs = self.pairs[:-max_pairs] if len(self.pairs) > max_pairs else self.pairs[:1]
        max_summary = CFG.get("max_summary_tokens", 600) * 4

        # Extract tool patterns for efficiency (v2)
        tool_patterns = set()
        blocks = []
        
        for u, a in old_pairs[:5]:
            # Detect tool calls in assistant output
            tool_markers = []
            for tool in ["web_search", "cached_web_search", "run_terminal", "read_file", "write_file"]:
                if f'"{tool}"' in a or tool + "(" in a:
                    tool_markers.append(tool)
            
            if tool_markers:
                tool_patterns.add(",".join(tool_markers))
            
            # Truncate user query
            u_trunc = u[:150].replace(chr(10), ' ')
            # Truncate assistant — keep key info
            a_trunc = a[:300].replace(chr(10), ' ')
            
            if tool_markers:
                a_trunc = f"[TOOLS:{','.join(tool_markers)}] {a_trunc}"
            
            blocks.append(f"User: {u_trunc}\nAssistant: {a_trunc}")

        block = "\n\n".join(blocks)

        # Add tool pattern hints at top (token-efficient recall)
        if tool_patterns:
            block = f"[TOOL HISTORY: {' | '.join(list(tool_patterns)[:3])}]\n\n{block}"

        if self.summary:
            self.summary = f"{self.summary}\n\n---\n{block}"
        else:
            self.summary = block

        # Truncate to max_summary
        if len(self.summary) > max_summary:
            self.summary = self.summary[:max_summary]

        # Update in-memory pairs
        self.pairs = self.pairs[-max_pairs:] if len(self.pairs) > max_pairs else self.pairs
        _db.update_summary(self.id, self.summary)
        self.save()
        return f"compressed {len(old_pairs)} old turns, {len(tool_patterns)} tool patterns preserved"

    def to_messages(self) -> list[ModelMessage]:
        """Convert session to pydantic-ai messages with memory injection."""
        from pydantic_ai.messages import ModelRequest, ModelResponse, UserPromptPart, TextPart
        msgs: list[ModelMessage] = []

        # Inject global MEMORY.md
        if MEMORY_FILE.exists():
            mem = MEMORY_FILE.read_text(encoding="utf-8", errors="ignore")[:3000]
            if mem:
                msgs.append(ModelRequest(parts=[UserPromptPart(content="[Persistent Memory]")]))
                msgs.append(ModelResponse(parts=[TextPart(content=mem)]))

        # Inject global USER.md
        if USER_FILE.exists():
            usr = USER_FILE.read_text(encoding="utf-8", errors="ignore")[:1500]
            if usr:
                msgs.append(ModelRequest(parts=[UserPromptPart(content="[User Profile]")]))
                msgs.append(ModelResponse(parts=[TextPart(content=usr)]))

        # Inject session-specific memory
        if self.memory_md:
            msgs.append(ModelRequest(parts=[UserPromptPart(content="[Session Memory]")]))
            msgs.append(ModelResponse(parts=[TextPart(content=self.memory_md[:3000])]))

        # Inject summary
        if self.summary:
            msgs.append(ModelRequest(parts=[UserPromptPart(content="[Prior Context Summary]")]))
            msgs.append(ModelResponse(parts=[TextPart(content=self.summary)]))

        # Recent pairs
        for u, a in self.pairs[-CFG["max_pairs"]:]:
            msgs.append(ModelRequest(parts=[UserPromptPart(content=u)]))
            msgs.append(ModelResponse(parts=[TextPart(content=a)]))
        return msgs

    def add(self, user: str, assistant: str):
        self.pairs.append((user, assistant))
        _db.add_pair(self.id, user, assistant)
        if len(self.pairs) > CFG["max_pairs"] + 2:
            self.compress()
        self.save()

    @property
    def token_est(self) -> int:
        """Estimate token count. Uses 3 chars/token for better English accuracy."""
        total = (len(self.summary) + len(self.memory_md) + len(self.user_md)
                 + sum(len(u) + len(a) for u, a in self.pairs))
        # Also check global memory files
        if MEMORY_FILE.exists():
            total += min(len(MEMORY_FILE.read_text(errors="ignore")), 3000)
        if USER_FILE.exists():
            total += min(len(USER_FILE.read_text(errors="ignore")), 1500)
        return total // 3

    @staticmethod
    def list_all() -> list[dict]:
        """Return all sessions from SQLite, sorted by last_active."""
        return _db.list_sessions()


# ══════════════════════════════════════════════════════════════════════════
#  MULTILINE INPUT  (raw mode + bracketed paste)
# ══════════════════════════════════════════════════════════════════════════

PASTE_START = b"\x1b[200~"
PASTE_END   = b"\x1b[201~"
BP_ON       = "\x1b[?2004h"
BP_OFF      = "\x1b[?2004l"


class InputBuffer:
    """
    Raw-mode terminal input with bracketed paste support.
    Pasted text is cleaned before it reaches the prompt so it stays editable.
    """

    def _redraw(self, prompt: str, buf: list[str]):
        _w(CLEAR_LINE)
        _w(prompt + "".join(buf))

    def readline(self, prompt_prefix: str = "") -> Optional[str]:
        if not sys.stdin.isatty():
            try:
                return input().strip() or None
            except EOFError:
                return None

        fd = sys.stdin.fileno()
        old = termios.tcgetattr(fd)
        prompt = f"{C_GREEN}{BOLD}You {RESET}{DIM}> {RESET}"

        buf: list[str] = []
        in_paste = False
        accum = bytearray()

        _w(prompt)
        _w(BP_ON)
        sys.stdout.flush()

        try:
            tty.setraw(fd)
            while True:
                if not select.select([sys.stdin], [], [], 0.05)[0]:
                    continue

                chunk = os.read(fd, 1024)
                if not chunk:
                    continue
                accum.extend(chunk)

                while accum:
                    if not in_paste and accum.startswith(PASTE_START):
                        del accum[:len(PASTE_START)]
                        in_paste = True
                        continue

                    if in_paste:
                        end_idx = accum.find(PASTE_END)
                        if end_idx == -1:
                            paste = bytes(accum)
                            accum.clear()
                            text = paste.decode("utf-8", errors="ignore")
                        else:
                            paste = bytes(accum[:end_idx])
                            del accum[:end_idx + len(PASTE_END)]
                            in_paste = False
                            text = paste.decode("utf-8", errors="ignore")

                        # Keep pasted content clean and editable as a single line.
                        text = re.sub(r"\s+", " ", text).strip()
                        if text:
                            buf.extend(text)
                            _w(text)
                        continue

                    b = accum[0]
                    del accum[0]

                    if b in (13, 10):
                        termios.tcsetattr(fd, termios.TCSADRAIN, old)
                        _w(BP_OFF + "\r\n")
                        return "".join(buf).strip()

                    if b == 4:
                        termios.tcsetattr(fd, termios.TCSADRAIN, old)
                        _w(BP_OFF + "\r\n")
                        return None

                    if b == 3:
                        raise KeyboardInterrupt

                    if b in (127, 8):
                        if buf:
                            buf.pop()
                            self._redraw(prompt, buf)
                        continue

                    if b == 27:
                        # Swallow ANSI escape sequences from arrows/navigation.
                        # Common arrow keys arrive as ESC [ A/B/C/D.
                        if len(accum) >= 2 and accum[0] == 91:
                            del accum[:2]
                        elif accum:
                            del accum[:1]
                        continue

                    if 32 <= b < 127:
                        ch = chr(b)
                        buf.append(ch)
                        _w(ch)

        except KeyboardInterrupt:
            termios.tcsetattr(fd, termios.TCSADRAIN, old)
            _w(BP_OFF + "\r\n")
            raise
        except Exception:
            termios.tcsetattr(fd, termios.TCSADRAIN, old)
            _w(BP_OFF + "\r\n")
            raise
        finally:
            try:
                termios.tcsetattr(fd, termios.TCSADRAIN, old)
            except Exception:
                pass

_ibuf = InputBuffer()


# ══════════════════════════════════════════════════════════════════════════
#  STREAMING OUTPUT
# ══════════════════════════════════════════════════════════════════════════

_FENCE_RE = re.compile(r"```(?P<lang>\w+)?(?:\s+(?P<file>[^\n]+))?\n(?P<body>.*?)```", re.DOTALL)


def _render_code_block(m) -> None:
    """Render a single fenced code block with syntax highlighting + truncation."""
    lang = (m.group("lang") or "text").lower()
    file = m.group("file") or ""
    body = m.group("body").rstrip()

    # Truncate very long blocks to avoid flooding the terminal
    lines = body.splitlines()
    MAX_LINES = 120
    if len(lines) > MAX_LINES:
        body = "\n".join(lines[:MAX_LINES])
        body += f"\n{DIM}  ... [{len(lines) - MAX_LINES} more lines truncated] ...{RESET}"

    hdr = Text()
    hdr.append(f"  {lang}", style="bold cyan")
    if file:
        hdr.append(f"  {file}", style="bold white")
    console.print(hdr)
    console.print(Syntax(
        body, lang,
        theme="monokai",
        line_numbers=len(lines) > 8,
        word_wrap=True,
        background_color="default",
    ))
    _ln()


def _render_response(text: str):
    """Render the full assistant response with proper markdown formatting."""
    text = _normalize_output(text)
    if not text:
        return

    global _was_streamed_live
    was_live = _was_streamed_live[0]
    if was_live:
        _was_streamed_live[0] = False  # reset for next turn

    blocks = list(_FENCE_RE.finditer(text))

    if was_live:
        # Raw text already printed in real-time — add syntax-highlighted code blocks
        if blocks:
            _ln()
            _ln(f"{DIM}── highlighted ──{RESET}")
            for m in blocks:
                _render_code_block(m)
        return

    # Non-streamed: render full markdown in a styled panel
    console.print()
    console.print(Panel(
        Markdown(text),
        title="[dim]Response[/dim]",
        border_style="border",
        box=box.ROUNDED,
        padding=(1, 2),
    ))

    if not blocks:
        return

    _ln()
    _ln(f"{DIM}── highlighted ──{RESET}")
    for m in blocks:
        _render_code_block(m)


async def _stream(query: str, history: list[ModelMessage]) -> tuple[str, bool]:
    _w(f"\n{C_CYAN}{BOLD}Agent{RESET}\n")

    full = ""

    # Patch tool badge for thinking logs
    _patch_td_with_thinking()

    try:
        with Live(
            Panel(
                Text("Thinking…", style="dim"),
                border_style="border",
                box=box.ROUNDED,
                padding=(1, 2),
            ),
            console=console,
            transient=True,
            refresh_per_second=18,
        ) as live:
            async with agent.run_stream(query, message_history=history) as result:
                async for chunk in result.stream_text():
                    full = chunk if chunk.startswith(full) else full + chunk

                    # Show thinking — full text, no 200-char cap
                    display_text = _normalize_output(full) or "Thinking…"
                    char_count = len(display_text)
                    # Truncate at word boundary if > 1200 chars to keep panel readable
                    if len(display_text) > 1200:
                        cut = display_text[:1200]
                        last_space = cut.rfind(' ')
                        display_text = cut[:last_space] + "…"
                        char_count = len(display_text)

                    live.update(Panel(
                        Text(display_text, style="dim"),
                        title=Text(f" streaming  ·  {char_count:,} chars ", style="dim"),
                        border_style="border",
                        box=box.ROUNDED,
                        padding=(1, 2),
                    ))

        return full, True

    except Exception as e:
        _w(f"{C_RED}Stream error: {e}{RESET}\n")
        return full, False

    finally:
        _unpatch_td()


# ══════════════════════════════════════════════════════════════════════════
#  PARALLEL STREAMING LOOP  (raw OpenAI SDK + ThreadPoolExecutor)
# ══════════════════════════════════════════════════════════════════════════

def _stream_parallel(query: str, history: list[ModelMessage],
                     system_prompt: str) -> tuple[str, bool]:
    """
    Stream response using the raw OpenAI SDK + parallel tool execution.
    Prints text in real-time via text_callback.

    UX: shows a thinking spinner until the first token arrives,
    then switches to a streaming counter. Both indicators run on the
    main thread (no separate spinner thread needed).
    """
    import concurrent.futures
    import queue
    import threading
    import time as time_module
    import sys
    import shutil

    from open_agent.agent_loop import AgentLoop, TOOL_SCHEMAS

    full = ""
    streamed_text = []  # Accumulate for markdown rendering
    streamed_len = [0]

    # Spinner state — tracks transition from "thinking" → "streaming"
    SPIN_CHARS   = "⠋⠙⠹⠸⠼⠴⠦⠧⠇⠏"
    THINK_PHRASES = ["thinking", "processing", "reasoning", "working on it",
                     "analyzing", "computing", "figuring it out"]
    spin_idx   = [0]
    phrase_idx = [0]
    last_phrase_change = [time_module.monotonic()]
    first_token  = threading.Event()   # set on first real text chunk
    done         = [False]

    # Reasoning box state (hermes-agent pattern)
    reasoning_queue: queue.Queue = queue.Queue()
    reasoning_buf   = ""   # accumulated raw reasoning content
    reasoning_box_open = False
    reasoning_box_closed_by_content = False  # once content starts, don't reopen

    def _on_text_chunk(chunk: str):
        if not chunk:
            return
        # Signal thinking spinner to retire on first real content
        if not first_token.is_set():
            first_token.set()
        _set_streamed()
        streamed_text.append(chunk)
        streamed_len[0] += len(chunk)
        # Real-time token output — direct stdout, no Rich buffering
        try:
            sys.stdout.write(chunk)
            sys.stdout.flush()
        except OSError:
            pass

    def _on_reasoning_chunk(text: str):
        """Called from agent thread — push reasoning text to main thread."""
        reasoning_queue.put(text)

    def _render_reasoning_box(text: str, opening: bool = False, closing: bool = False):
        """Render a reasoning box line — hermes-agent-style dim thinking box."""
        try:
            w = shutil.get_terminal_size().columns
        except Exception:
            w = 80
        prefix = "  "
        r_label = " Thinking "
        r_width = max(w - len(prefix) - 4, 10)
        if opening:
            # Top border: "  ┌─ Thinking ─────────────┐"
            dash_count = max(r_width - len(r_label) - 1, 1)
            _w(f"\n{DIM}{prefix}┌─{r_label}{'─' * dash_count}┐{RESET}\n")
        if closing:
            dash_count = max(r_width - len(r_label) - 1, 1)
            _w(f"\n{DIM}{prefix}└{'─' * (len(r_label) + dash_count + 1)}┘{RESET}\n")
        # Content line
        for line in text.splitlines():
            if line.strip():
                _w(f"{DIM}{prefix}│ {line[:w - len(prefix) - 4]}{' ' * max(0, w - len(prefix) - 4 - len(line[:w - len(prefix) - 4]))}{DIM}│{RESET}\n")
        if text.strip() and not text.splitlines():
            _w(f"{DIM}{prefix}│ {text.strip()[:w - len(prefix) - 4]}{DIM}│{RESET}\n")

    loop = AgentLoop(
        base_url=CFG["base_url"],
        api_key="lm-studio",
        model=CFG["model_name"],
        system_prompt=system_prompt,
        tools=TOOL_SCHEMAS,
        max_tokens=2400,
        temperature=CFG.get("agent_temperature", 0.6),
        text_callback=_on_text_chunk,
        reasoning_callback=_on_reasoning_chunk,
    )

    result_container = [None]
    exc_container   = [None]

    def run_loop():
        try:
            result_container[0] = loop.run(query, history)
        except Exception as e:
            exc_container[0] = e
        finally:
            done[0] = True

    thread = threading.Thread(target=run_loop, name="agent-loop")
    thread.start()

    # ── Unified spinner + reasoning box: main thread ──────────────────────
    CLEAR_EOL = "\x1b[K"
    try:
        while thread.is_alive() or not first_token.is_set():
            now = time_module.monotonic()

            # Drain reasoning queue (from agent thread)
            while not reasoning_queue.empty():
                chunk = reasoning_queue.get_nowait()
                if not reasoning_box_open and not reasoning_box_closed_by_content:
                    # Open the thinking box
                    reasoning_box_open = True
                    _render_reasoning_box("", opening=True)
                if reasoning_box_open:
                    # Accumulate and display each line
                    reasoning_buf += chunk
                    # Flush complete lines to display
                    while "\n" in reasoning_buf:
                        line, reasoning_buf = reasoning_buf.split("\n", 1)
                        if line.strip():
                            _render_reasoning_box(line)
                    # Flush partial lines > 80 chars for real-time visibility
                    if len(reasoning_buf) > 80:
                        _render_reasoning_box(reasoning_buf)
                        reasoning_buf = ""

            # Once first text token arrives, close the reasoning box
            if first_token.is_set() and reasoning_box_open and not reasoning_box_closed_by_content:
                # Flush remaining buf
                if reasoning_buf.strip():
                    _render_reasoning_box(reasoning_buf)
                _render_reasoning_box("", closing=True)
                reasoning_box_closed_by_content = True
                reasoning_box_open = False

            ch = SPIN_CHARS[spin_idx[0] % len(SPIN_CHARS)]

            if not first_token.is_set():
                # Still waiting for first token — show thinking indicator
                if now - last_phrase_change[0] > 3.0:
                    phrase_idx[0] = (phrase_idx[0] + 1) % len(THINK_PHRASES)
                    last_phrase_change[0] = now
                phrase = THINK_PHRASES[phrase_idx[0]]
                label = f"  {ch}  {phrase}..."
            else:
                # First token received — show streaming counter
                chars = streamed_len[0]
                label = f"  {ch}  streaming  ·  {chars:,} chars"

            _w(f"\r{DIM}{label}{CLEAR_EOL}{RESET}   ")
            spin_idx[0] += 1
            time_module.sleep(0.08)
    finally:
        done[0] = True
        # Close reasoning box if still open when loop exits
        if reasoning_box_open:
            if reasoning_buf.strip():
                _render_reasoning_box(reasoning_buf)
            _render_reasoning_box("", closing=True)

    thread.join()
    exc = exc_container[0]

    # Clear the spinner line
    _w(f"\r{' ' * 70}\r")

    if exc:
        _w(f"{C_RED}Parallel loop error: {exc}{RESET}\n")
        return "", False

    full = result_container[0] or "".join(streamed_text)
    return full, bool(full)

class AgentSession:
    def __init__(self, session: Session):
        self.session = session
        self._turns_since_memory_nudge = 0

    def _augment(self, query: str) -> str:
        q = query.lower()
        hints = []
        if any(t in q for t in GROUNDING_TRIGGERS):
            hints.append("[Grounding required: run date via terminal first, then search.]")
            hints.append("[Execution mode: use grounded facts directly. Do not narrate plans.]")
        if any(t in q for t in SOUL_TRIGGERS):
            hints.append("[Complex task: call load_soul before proceeding.]")
            hints.append("[Execution mode: tool-first, concise, action-oriented.]")

        # Memory nudge — prompt model to review memory every N turns
        nudge_interval = CFG.get("memory_nudge_interval", 10)
        self._turns_since_memory_nudge += 1
        if self._turns_since_memory_nudge >= nudge_interval:
            hints.append("[Memory nudge: Consider reviewing or updating MEMORY.md with any new facts from this conversation.]")
            self._turns_since_memory_nudge = 0

        return query + ("\n\n" + "\n".join(hints) if hints else "")

    async def run(self, query: str) -> Optional[str]:
        # Clear any stale interrupt flag before starting a new turn
        from open_agent.agent_loop import clear_interrupt
        clear_interrupt()
        _ln(); _w(f"{C_SLATE}{'─' * 60}{RESET}\n")
        sess = self.session
        tok = sess.token_est + len(query) // 4

        if tok > CFG["ctx_compress"]:
            td.info(f"context guard: {sess.compress()}")
        elif tok > CFG["ctx_warn"]:
            td.info(f"context ~{tok:,} tokens — approaching limit")

        _w(f"  {DIM}session {C_PURPLE}{sess.id}{RESET}\n")

        # Show session memory status
        if sess.memory_md:
            td.info(f"session memory: {len(sess.memory_md)} chars")
        if sess.summary:
            td.info(f"summary active: {len(sess.summary)} chars")

        augmented = self._augment(query)
        history = sess.to_messages()

        # Proactive grounding for time-sensitive/current queries.
        grounding_blob = ""
        if any(t in query.lower() for t in GROUNDING_TRIGGERS):
            td.info("execution mode: grounding tools")
            try:
                date_output = await run_terminal(None, "date && uname -r")
                search_output = await web_search(None, query)
                grounding_blob = (
                    "[Grounded snapshot]\n"
                    f"{date_output}\n\n"
                    "[Search snapshot]\n"
                    f"{search_output}"
                )
                augmented = (
                    f"{augmented}\n\n"
                    f"{grounding_blob}\n\n"
                    "Use the grounded snapshot above directly. Do not narrate what you will do."
                )
            except Exception as e:
                td.info(f"grounding bootstrap failed: {e}")

        # Complex tasks can load the extended behavioral layer before the model thinks.
        if any(t in query.lower() for t in SOUL_TRIGGERS):
            try:
                td.info("loading extended instructions")
                soul_text = await load_soul(None)
                augmented = (
                    f"{augmented}\n\n"
                    "[Extended instructions]\n"
                    f"{soul_text}\n\n"
                    "Follow the execution rules above and answer directly."
                )
            except Exception as e:
                td.info(f"soul bootstrap failed: {e}")

        for attempt in range(3):
            use_parallel = CFG.get("use_parallel_loop", False)
            if use_parallel:
                # Parallel loop: runs sync but we call from async context
                import asyncio
                loop = asyncio.get_event_loop()
                full, ok = await loop.run_in_executor(
                    None, lambda: _stream_parallel(augmented, history, SYSTEM_PROMPT)
                )
            else:
                full, ok = await _stream(augmented, history)

            if ok and full:
                full = _normalize_output(full)

                # Fallback: if the model starts planning, strengthen the prompt and rerun once.
                if _looks_like_planning(full):
                    td.info("model is planning — forcing execution")
                    try:
                        forced_augmented = (
                            "EXECUTION MODE ACTIVE. "
                            "Do not explain what you will do. Answer directly using the grounded information and tools."
                        )
                        if grounding_blob:
                            forced_augmented += f"\n\n{grounding_blob}"
                        forced_augmented += f"\n\nUser request: {query}"
                        if use_parallel:
                            loop2 = asyncio.get_event_loop()
                            full, ok = await loop2.run_in_executor(
                                None, lambda: _stream_parallel(forced_augmented, history, SYSTEM_PROMPT)
                            )
                        else:
                            full, ok = await _stream(forced_augmented, history)
                        full = _normalize_output(full)
                    except Exception as e:
                        td.info(f"forced execution failed: {e}")

                if ok and full:
                    sess.add(query, full)
                    _render_response(full)
                    return full

            if not ok and attempt < 2:
                _w(f"\n{DIM}Retry {attempt+1}/3…{RESET}\n")
                await asyncio.sleep(1.5)

        _w(f"{C_RED}✗ Failed after 3 attempts.{RESET}")
        return None


# ══════════════════════════════════════════════════════════════════════════
#  FIRST-RUN ONBOARDING WIZARD
# ══════════════════════════════════════════════════════════════════════════

def _ask(prompt: str, default: str = "") -> str:
    """Simple interactive prompt with a default value."""
    hint = f" [{default}]" if default else ""
    _w(f"  {C_CYAN}{prompt}{hint}{RESET}: ")
    try:
        val = input().strip()
        return val or default
    except (EOFError, KeyboardInterrupt):
        return default

def _yn(prompt: str, default: bool = True) -> bool:
    hint = "Y/n" if default else "y/N"
    _w(f"  {C_CYAN}{prompt}{RESET} [{hint}]: ")
    try:
        val = input().strip().lower()
        if not val: return default
        return val in ("y", "yes")
    except (EOFError, KeyboardInterrupt):
        return default

def run_setup_wizard(cfg: dict, first_run: bool = True) -> dict:
    """Interactive first-run configuration wizard."""
    _ln()
    if first_run:
        console.print(Panel(
            f"[header]Welcome to Open-Agent![/header]\n\n"
            f"[dim]Let's get you set up in about 60 seconds.\n"
            f"Press Enter to keep defaults. You can change anything later with [/dim][tool.name]/setup[/tool.name]",
            box=box.ROUNDED, border_style="border", padding=(1, 3),
        ))
    else:
        console.print(Panel(
            "[header]Open-Agent Setup[/header]\n[dim]Reconfigure your settings.[/dim]",
            box=box.ROUNDED, border_style="border", padding=(1, 2),
        ))

    _ln()

    # ── Step 1: Inference server ──────────────────────────────────────────
    _w(f"\n{C_AMBER}{BOLD}[1/3] Inference Server{RESET}\n")
    _w(f"  {DIM}Open-Agent connects to any OpenAI-compatible API.\n")
    _w(f"  Default: llama.cpp on localhost:8083{RESET}\n\n")

    base_url = _ask("Server URL", cfg.get("base_url", DEFAULT_CONFIG["base_url"]))
    model_name = _ask("Model name", cfg.get("model_name", DEFAULT_CONFIG["model_name"]))

    # Quick connectivity check
    _w(f"\n  {DIM}Checking connection…{RESET}")
    try:
        import urllib.request
        urllib.request.urlopen(base_url.rstrip("/") + "/models", timeout=3)
        _w(f"  {C_GREEN}✓ Connected{RESET}\n")
    except Exception:
        _w(f"  {C_AMBER}⚠ Could not reach server (make sure llama.cpp is running){RESET}\n")

    # ── Step 2: Web search ────────────────────────────────────────────────
    _w(f"\n{C_AMBER}{BOLD}[2/3] Web Search (SearxNG){RESET}\n")
    _w(f"  {DIM}SearxNG is a private, self-hosted search engine.\n")
    _w(f"  Run it with: docker run -d -p 8081:8080 searxng/searxng{RESET}\n\n")

    searxng_url = _ask("SearxNG URL", cfg.get("searxng_url", DEFAULT_CONFIG["searxng_url"]))

    # ── Step 3: Obsidian (optional) ───────────────────────────────────────
    _w(f"\n{C_AMBER}{BOLD}[3/3] Obsidian Vault (optional){RESET}\n")
    _w(f"  {DIM}Connect your Obsidian vault so the agent can read and write your notes.{RESET}\n\n")

    use_obsidian = _yn("Connect Obsidian vault?", default=bool(cfg.get("obsidian_path")))
    obsidian_path = ""
    if use_obsidian:
        obsidian_path = _ask("Vault path", cfg.get("obsidian_path", str(Path.home() / "Documents" / "Obsidian")))
        if obsidian_path and not Path(obsidian_path).exists():
            _w(f"  {C_AMBER}⚠ Path doesn't exist yet — you can create it later{RESET}\n")
        elif obsidian_path:
            _w(f"  {C_GREEN}✓ Vault found{RESET}\n")

    # ── Save ──────────────────────────────────────────────────────────────
    cfg.update({
        "base_url":     base_url,
        "model_name":   model_name,
        "searxng_url":  searxng_url,
        "obsidian_path":obsidian_path,
        "setup_done":   True,
    })
    save_config(cfg)

    _ln()
    console.print(Panel(
        f"[tool.ok]✓ Setup complete![/tool.ok]\n\n"
        f"[dim]Config saved to[/dim] [accent]{CONFIG_FILE}[/accent]\n"
        f"[dim]Run [/dim][tool.name]/setup[/tool.name][dim] at any time to change settings.\n"
        f"[dim]Type[/dim] [tool.name]/help[/tool.name] [dim]to see all commands.[/dim]",
        box=box.ROUNDED, border_style="border", padding=(1, 3),
    ))
    _ln()
    return cfg


# ══════════════════════════════════════════════════════════════════════════
#  BANNER
# ══════════════════════════════════════════════════════════════════════════

VERSION = "1.1"

LOGO_WIDE = r"""
┌───────────────────────────────────────────────────────────────┐
│                                                               │
│   ░░▒▒▓▓  █▀█ █▀█ █▀▀ █▄ █    ▄▀█ █▀▀ █▀▀ █▄ █ ▀█▀   ▓▓▒▒░░  │
│   ░░▒▒▓▓  █▄█ █▀▀ ██▄ █ ▀█    █▀█ █▄█ ██▄ █ ▀█  █    ▓▓▒▒░░  │
│                                                               │
│       local-first · privacy-first · intelligence-driven       │
│                                                               │
└───────────────────────────────────────────────────────────────┘
"""

LOGO_NARROW = r"""
┌─────────────────────────────────────────────────────┐
│  ░▒▓ █▀█ █▀█ █▀▀ █▄ █ ▄▄ ▄▀█ █▀▀ █▀▀ █▄ █ ▀█▀ ▓▒░  │
│  ░▒▓ █▄█ █▀▀ ██▄ █ ▀█    █▀█ █▄█ ██▄ █ ▀█  █  ▓▒░  │
└─────────────────────────────────────────────────────┘
"""

# ── ANSI constants ────────────────────────────────────────────────
BRIGHT = '\033[1;97m'   # bold bright white — glowing banner
DIM    = '\033[2m'      # dark grey         — tagline / labels
RESET  = '\033[0m'


def _terminal_width() -> int:
    try:
        return shutil.get_terminal_size().columns
    except Exception:
        return 80


def print_banner(session_id: str) -> None:
    width = _terminal_width()

    _ln()

    if width >= 100:
        for line in LOGO_WIDE.splitlines():
            pad = max(0, (width - len(line)) // 2)
            _w(f"{BRIGHT}{' ' * pad}{line}{RESET}\n")
    else:
        lines = LOGO_NARROW.splitlines()
        mid = len(lines) // 2
        for i, line in enumerate(lines):
            pad = max(0, (width - len(line)) // 2)
            color = BRIGHT if i <= mid else DIM
            _w(f"{color}{' ' * pad}{line}{RESET}\n")

    # Tagline
    tagline_left  = "Run capable AI on any laptop — "
    tagline_right = (
        f"{BRIGHT}private{RESET}{DIM}, {RESET}"
        f"{BRIGHT}free{RESET}{DIM}, {RESET}"
        f"{BRIGHT}yours{RESET}{DIM}.{RESET}"
    )
    tagline_visible_len = len(tagline_left) + len("private, free, yours.")
    pad = max(0, (width - tagline_visible_len) // 2)
    _w(f"\n{' ' * pad}{DIM}{tagline_left}{RESET}{tagline_right}\n\n")

    # Session ID
    _w(f"  {DIM}Session{RESET}  {BRIGHT}{session_id}{RESET}\n")

    # Commands
    _w(f"\n  {DIM}Type{RESET} ")
    _w(f"{BRIGHT}/help{RESET}")
    _w(f"  {DIM}·{RESET}  ")
    _w(f"{BRIGHT}/sessions{RESET}")
    _w(f"  {DIM}·{RESET}  ")
    _w(f"{BRIGHT}/exit{RESET}")
    _w(f"  {DIM}to quit{RESET}\n\n")


# ══════════════════════════════════════════════════════════════════════════
#  SLASH COMMANDS
# ══════════════════════════════════════════════════════════════════════════

def _help_table() -> Table:
    tbl = Table(box=box.SIMPLE, show_header=True, padding=(0, 2), header_style="header")
    tbl.add_column("Command",     style="tool.name", no_wrap=True)
    tbl.add_column("Description", style="dim")
    tbl.add_column("Alias",       style="dim")
    for cmd, desc, alias in [
        ("/help",     "Show this help",              "/h"),
        ("/sessions", "Browse and resume old chats", "/s"),
        ("/new",      "Start a fresh session",       ""),
        ("/history",  "Token count + session info",  "/hi"),
        ("/compress", "Force compress history",       ""),
        ("/clear",    "Wipe current session",         "/cl"),
        ("/soul",     "Show SOUL.md",                 ""),
        ("/setup",    "Reconfigure settings",         ""),
        ("/model",    "Show model + server config",   "/m"),
        ("/sources",  "List RSS sources",             ""),
        ("/save",     "Force save current session",   ""),
        ("/exit",     "Save and quit",                "/q"),
    ]:
        tbl.add_row(cmd, desc, alias)
    return tbl


def _sessions_panel(runner: "AgentRunner") -> bool:
    """Show session browser and allow resuming an old session."""
    sessions = Session.list_all()
    if not sessions:
        console.print("[dim]No sessions found.[/dim]")
        return True

    tbl = Table(box=box.SIMPLE, show_header=True, padding=(0, 2), header_style="header")
    tbl.add_column("#",           style="dim", no_wrap=True)
    tbl.add_column("Session ID",  style="tool.name", no_wrap=True)
    tbl.add_column("Title",       style="dim")
    tbl.add_column("Turns",       style="dim", no_wrap=True)
    tbl.add_column("Last active", style="dim", no_wrap=True)

    for i, s in enumerate(sessions[:15], 1):
        try:
            la = datetime.fromisoformat(s["last_active"]).strftime("%b %d %H:%M")
        except Exception:
            la = s["last_active"][:16]
        current = " ← current" if s["id"] == runner.session.id else ""
        tbl.add_row(
            str(i),
            s["id"] + current,
            s["title"][:45] or "(no title)",
            str(s["turns"]),
            la,
        )

    console.print(Panel(tbl, title="[header]Sessions[/header]",
                        border_style="border", box=box.ROUNDED))
    _w(f"\n  {DIM}Enter number to resume, or press Enter to stay in current session:{RESET} ")
    try:
        val = input().strip()
        if val.isdigit():
            idx = int(val) - 1
            if 0 <= idx < len(sessions):
                sid = sessions[idx]["id"]
                if sid != runner.session.id:
                    runner.session.save()
                    runner.session = Session(sid)
                    _w(f"  {C_GREEN}✓ Resumed session {sid}{RESET}\n")
                    td.info(f"loaded {len(runner.session.pairs)} turns")
                else:
                    _w(f"  {DIM}Already in that session.{RESET}\n")
    except (EOFError, KeyboardInterrupt):
        pass
    return True


async def handle_slash(cmd: str, runner: "AgentRunner") -> bool:
    """Handle slash commands (/help, /setup, /exit, etc.)."""
    global CFG   # Needed because we reassign CFG in /setup

    parts = cmd.strip().split()
    c = parts[0].lower()

    if c in ("/exit", "/quit", "/q"):
        runner.session.save()
        console.print(Panel("[dim]Session saved. Goodbye.[/dim]", border_style="border"))
        sys.exit(0)

    if c in ("/help", "/h"):
        console.print(Panel(_help_table(), title="[header]Open-Agent[/header]",
                            border_style="border", box=box.ROUNDED))
        return True

    if c in ("/sessions", "/s"):
        return _sessions_panel(runner)

    if c == "/new":
        runner.session.save()
        runner.session = Session()
        _w(f" {C_GREEN}✓ New session started: {runner.session.id}{RESET}\n")
        return True

    if c in ("/clear", "/cl"):
        runner.session.pairs.clear()
        runner.session.summary = ""
        runner.session.save()
        _w(f" {C_GREEN}✓ Session cleared.{RESET}\n")
        return True

    if c == "/compress":
        msg = runner.session.compress()
        _w(f" {C_GREEN}✓ {msg}{RESET}\n")
        return True

    if c in ("/history", "/hi"):
        s = runner.session
        tok = s.token_est
        tbl = Table(box=box.SIMPLE, show_header=False, padding=(0, 2))
        tbl.add_column(style="tool.name")
        tbl.add_column(style="dim")
        tbl.add_row("Session ID", s.id)
        tbl.add_row("Title", s.title or "(none yet)")
        tbl.add_row("Verbatim turns", str(len(s.pairs)))
        tbl.add_row("Summary", "yes" if s.summary else "no")
        tbl.add_row("Token estimate", f"~{tok:,}")
        tbl.add_row("Summary size", f"{len(s.summary):,} chars")
        tbl.add_row("Context used", f"{tok / CFG['ctx_limit'] * 100:.1f}%")
        tbl.add_row("Context limit", f"{CFG['ctx_limit']:,}")
        console.print(Panel(tbl, title="[header]History[/header]",
                            border_style="border", box=box.ROUNDED))
        return True

    if c == "/soul":
        if SOUL_FILE.exists():
            console.print(Panel(Markdown(SOUL_FILE.read_text()),
                                title="[header]SOUL.md[/header]",
                                border_style="border", box=box.ROUNDED, padding=(1, 2)))
        else:
            _w(f" {C_RED}SOUL.md not found{RESET}\n")
        return True

    if c == "/setup":
        CFG = run_setup_wizard(CFG, first_run=False)
        return True

    if c in ("/model", "/m"):
        tbl = Table(box=box.SIMPLE, show_header=False, padding=(0, 2))
        tbl.add_column(style="tool.name")
        tbl.add_column(style="accent")
        for k, v in [("Model", CFG["model_name"]), ("Server", CFG["base_url"]),
                     ("Search", CFG["searxng_url"]), ("Context", f"{CFG['ctx_limit']:,} tokens"),
                     ("KV cache", "Q4_0 (K+V)"), ("Flash attn", "on")]:
            tbl.add_row(k, v)
        console.print(Panel(tbl, title="[header]Config[/header]",
                            border_style="border", box=box.ROUNDED))
        return True

    if c == "/sources":
        names = sorted(RSS_SOURCES.keys())
        cols = 4
        chunks = [names[i:i+cols] for i in range(0, len(names), cols)]
        tbl = Table(box=box.SIMPLE, show_header=False, padding=(0, 2))
        for _ in range(cols):
            tbl.add_column(style="tool.name")
        for chunk in chunks:
            tbl.add_row(*chunk + [""] * (cols - len(chunk)))
        cats = " · ".join(sorted(BROAD_CATS))
        console.print(Panel(tbl, title="[header]RSS Sources[/header]",
                            subtitle=f"[dim]categories: {cats}[/dim]",
                            border_style="border", box=box.ROUNDED))
        return True

    if c == "/save":
        runner.session.save()
        _w(f" {C_GREEN}✓ Saved session {runner.session.id}{RESET}\n")
        return True

    if c in ("/parallel", "/pll"):
        enabled = CFG.get("use_parallel_loop", False)
        enabled = not enabled
        CFG["use_parallel_loop"] = enabled
        save_config(CFG)
        status = f" {C_GREEN}ENABLED{RESET}" if enabled else f" {C_SLATE}disabled{RESET}"
        _w(f" Parallel tool mode:{status}\n")
        _w(f" {DIM}Uses raw OpenAI SDK + ThreadPoolExecutor for parallel tool calls.{RESET}\n")
        return True

    return False

# ══════════════════════════════════════════════════════════════════════════
#  RUNNER  (ties it all together)
# ══════════════════════════════════════════════════════════════════════════

class AgentRunner:
    def __init__(self, session: Session):
        self.session    = session
        self._agent_run = AgentSession(session)

    @property
    def _agent(self) -> AgentSession:
        # Always returns an AgentSession bound to the current session
        if self._agent_run.session is not self.session:
            self._agent_run = AgentSession(self.session)
        return self._agent_run


# ══════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════

async def _main_loop(runner: AgentRunner):
    while True:
        try:
            _ln()
            user_input = _ibuf.readline()
            if user_input is None:
                runner.session.save(); break
            if not user_input.strip(): continue
            if user_input.startswith("/"):
                if await handle_slash(user_input, runner): continue
            await runner._agent.run(user_input)
        except KeyboardInterrupt:
            from open_agent.agent_loop import interrupt_agent
            interrupt_agent()
            _w(f"\n{DIM}Ctrl-C — saving…{RESET}\n")
            runner.session.save()
            _w(f"{DIM}Goodbye.{RESET}\n"); break
        except EOFError:
            runner.session.save(); break
        except Exception as e:
            _w(f"{C_RED}Error: {e}{RESET}\n")


def main():
    global CFG
    parser = argparse.ArgumentParser(
        prog="opagent",
        description="Open-Agent — local AI agent for consumer hardware",
    )
    parser.add_argument("--setup",   action="store_true", help="Run setup wizard")
    parser.add_argument("--version", action="store_true", help="Print version")
    parser.add_argument("--new",     action="store_true", help="Start a new session")
    parser.add_argument("--resume",  metavar="SESSION_ID", help="Resume a specific session")
    parser.add_argument("--tui",     action="store_true", help="Launch the TUI mode")
    args = parser.parse_args()

    if args.version:
        print("open-agent 5.0.0"); sys.exit(0)
    
    # Launch TUI mode
    if args.tui:
        from open_agent.tui import run_tui
        run_tui()
        sys.exit(0)

    # First-run onboarding
    if not CFG.get("setup_done") or args.setup:
        CFG = run_setup_wizard(CFG, first_run=not CFG.get("setup_done"))

    # Session selection
    if args.resume:
        session = Session(args.resume)
        if not session.path.exists():
            _w(f"{C_RED}Session '{args.resume}' not found.{RESET}\n"); sys.exit(1)
    elif args.new:
        session = Session()
    else:
        # Resume most recent session if it exists, else create new
        all_sessions = Session.list_all()
        if all_sessions and not args.new:
            session = Session(all_sessions[0]["id"])
        else:
            session = Session()

    runner = AgentRunner(session)
    print_banner(session.id)
    if session.pairs:
        td.info(f"resumed session {session.id} — {len(session.pairs)} turns")

    asyncio.run(_main_loop(runner))
    _db.close()


def main_sync():
    """Entry point for console_scripts."""
    main()


if __name__ == "__main__":
    main()
